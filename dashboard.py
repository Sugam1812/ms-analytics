import streamlit as st
import streamlit.components.v1 as components
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import sqlite3
import os
import io
from datetime import datetime
import numpy as np

# ── PPT generation ────────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    _PPTX_OK = True
except ImportError:
    _PPTX_OK = False


def _rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_rect(slide, x, y, w, h, fill_hex=None, line_hex=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill_hex:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill_hex)
    else:
        shape.fill.background()
    if line_hex and line_w > 0:
        shape.line.color.rgb = _rgb(line_hex)
    else:
        shape.line.fill.background()
    return shape


def _add_text(slide, text, x, y, w, h, size=12, bold=False, color="#0F172A",
              align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = "Calibri"
    return tb


def generate_ppt(orders_df, items_df):
    """Generate a branded MS Enterprises PPT matching dashboard design."""
    if not _PPTX_OK:
        return None

    # ── Palette ───────────────────────────────────────────────
    C_BG     = "#F8FAFC"
    C_DARK   = "#0F172A"
    C_INDIGO = "#6366F1"
    C_CYAN   = "#22D3EE"
    C_AMBER  = "#F59E0B"
    C_GREEN  = "#10B981"
    C_PINK   = "#EC4899"
    C_SLATE  = "#334155"
    C_MUTED  = "#64748B"
    C_CARD   = "#FFFFFF"
    C_BORDER = "#E2E8F0"

    ACCENTS = [C_INDIGO, C_CYAN, C_AMBER, C_GREEN, C_PINK, "#F97316", "#A78BFA", "#EF4444"]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # completely blank

    # ══════════════════════════════════════════════════════════
    # SLIDE 1 — TITLE
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)

    # dark indigo background
    _add_rect(slide, 0, 0, 13.33, 7.5, fill_hex=C_DARK)
    # accent stripe left
    _add_rect(slide, 0, 0, 0.55, 7.5, fill_hex=C_INDIGO)
    # subtle gradient block
    _add_rect(slide, 0.55, 3.2, 12.78, 0.06, fill_hex=C_INDIGO)

    _add_text(slide, "⚡ MS ENTERPRISES", 1.1, 1.0, 10, 0.9,
              size=13, bold=True, color=C_CYAN, align=PP_ALIGN.LEFT)
    _add_text(slide, "Purchase Order Intelligence", 1.1, 1.8, 10, 1.2,
              size=36, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT)
    _add_text(slide, "Material Analytics Dashboard · All Financal Years",
              1.1, 3.0, 10, 0.5, size=14, color="#94A3B8", align=PP_ALIGN.LEFT)

    # KPI preview row at bottom
    total_rev = orders_df["total_amount"].sum()
    total_qty = items_df["quantity"].sum()
    n_cust    = orders_df["customer_name"].nunique()
    n_states  = orders_df["consignee_state"].nunique()
    n_orders  = len(orders_df)

    kpi_items = [
        ("₹" + (f"{total_rev/1e7:.1f} Cr" if total_rev >= 1e7 else f"{total_rev/1e5:.1f} L"),
         "Total Revenue", C_INDIGO),
        (f"{total_qty:,.0f} MT", "Total Qty", C_AMBER),
        (str(n_orders), "Total Orders", C_CYAN),
        (str(n_cust), "Factories", C_GREEN),
        (str(n_states), "States", C_PINK),
    ]
    kx = 1.1
    for val, lbl, col in kpi_items:
        _add_rect(slide, kx, 4.1, 2.1, 1.3, fill_hex="#1E2A45", line_hex=col, line_w=Pt(1.5))
        _add_text(slide, val, kx + 0.12, 4.2, 1.9, 0.55,
                  size=17, bold=True, color=col, align=PP_ALIGN.LEFT)
        _add_text(slide, lbl, kx + 0.12, 4.72, 1.9, 0.4,
                  size=10, color="#94A3B8", align=PP_ALIGN.LEFT)
        kx += 2.28

    _add_text(slide, f"Generated {datetime.now().strftime('%d %b %Y')}", 1.1, 6.8, 10, 0.4,
              size=9, color="#475569", align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════
    # SLIDE 2 — EXECUTIVE KPIs
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 7.5, fill_hex=C_BG)
    _add_rect(slide, 0, 0, 13.33, 0.75, fill_hex=C_INDIGO)
    _add_text(slide, "Executive KPIs", 0.3, 0.1, 9, 0.55,
              size=20, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT)
    _add_text(slide, "Key performance indicators across all financial years",
              0.3, 0.45, 9, 0.35, size=10, color="#A5B4FC", align=PP_ALIGN.LEFT)

    # FY revenue table
    fy_rev = orders_df.groupby("financial_year")["total_amount"].sum().sort_index()
    fy_qty = items_df.groupby("financial_year")["quantity"].sum().sort_index() if not items_df.empty else pd.Series()
    fy_ord = orders_df.groupby("financial_year").size().sort_index()

    # KPI cards row 1
    cards_r1 = [
        ("💰 Total Revenue",
         "₹" + (f"{total_rev/1e7:.2f} Cr" if total_rev >= 1e7 else f"{total_rev/1e5:.1f} L"),
         C_INDIGO),
        ("⚖️ Total Qty (MT)", f"{total_qty:,.1f} MT", C_AMBER),
        ("📦 Total Orders",   str(n_orders),           C_CYAN),
        ("🏭 Factories",      str(n_cust),             C_GREEN),
        ("🗺️ States",        str(n_states),           C_PINK),
    ]
    kx = 0.25
    for title, val, col in cards_r1:
        _add_rect(slide, kx, 1.0, 2.42, 1.5, fill_hex=C_CARD, line_hex=C_BORDER, line_w=Pt(1))
        _add_rect(slide, kx, 1.0, 2.42, 0.18, fill_hex=col)
        _add_text(slide, title, kx + 0.15, 1.22, 2.15, 0.4,
                  size=10, bold=True, color=C_MUTED, align=PP_ALIGN.LEFT)
        _add_text(slide, val, kx + 0.15, 1.6, 2.15, 0.7,
                  size=20, bold=True, color=C_DARK, align=PP_ALIGN.LEFT)
        kx += 2.57

    # FY breakdown table
    _add_text(slide, "Revenue by Financial Year", 0.25, 2.8, 8, 0.4,
              size=13, bold=True, color=C_DARK)

    col_labels = ["Financial Year", "Revenue (₹)", "Qty (MT)", "Orders"]
    col_xs     = [0.25, 3.5, 7.0, 10.0]
    col_ws     = [3.0, 3.2, 2.6, 2.5]

    # header row
    _add_rect(slide, 0.25, 3.3, 12.85, 0.38, fill_hex=C_INDIGO)
    for lbl, cx in zip(col_labels, col_xs):
        _add_text(slide, lbl, cx + 0.1, 3.33, 2.8, 0.35,
                  size=10, bold=True, color="#FFFFFF", align=PP_ALIGN.LEFT)

    row_y = 3.68
    for i, fy in enumerate(fy_rev.index):
        bg = C_CARD if i % 2 == 0 else "#F1F5F9"
        _add_rect(slide, 0.25, row_y, 12.85, 0.37, fill_hex=bg)
        rev = fy_rev[fy]
        qty = fy_qty.get(fy, 0) if not fy_qty.empty else 0
        ord_ = fy_ord.get(fy, 0)
        row_vals = [
            str(fy),
            "₹" + (f"{rev/1e7:.2f} Cr" if rev >= 1e7 else f"{rev/1e5:.1f} L"),
            f"{qty:,.1f}",
            str(ord_),
        ]
        for val, cx in zip(row_vals, col_xs):
            _add_text(slide, val, cx + 0.1, row_y + 0.04, 2.8, 0.3,
                      size=10, color=C_SLATE, align=PP_ALIGN.LEFT)
        row_y += 0.37

    # ══════════════════════════════════════════════════════════
    # SLIDE 3 — TOP MATERIALS
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 7.5, fill_hex=C_BG)
    _add_rect(slide, 0, 0, 13.33, 0.75, fill_hex=C_AMBER)
    _add_text(slide, "Top Materials by Revenue", 0.3, 0.1, 9, 0.55,
              size=20, bold=True, color="#FFFFFF")
    _add_text(slide, "Horizontal bar chart — sorted by total revenue",
              0.3, 0.45, 9, 0.35, size=10, color="#FEF3C7")

    if not items_df.empty:
        mat_rev = (items_df.groupby("material_name")["amount"]
                   .sum().sort_values(ascending=False).head(12))
        max_rev = mat_rev.max() if mat_rev.max() > 0 else 1
        bar_area_w = 8.5  # max bar width in inches
        bar_h      = 0.38
        bar_gap    = 0.12
        start_y    = 1.05

        for i, (mat, rev) in enumerate(mat_rev.items()):
            y = start_y + i * (bar_h + bar_gap)
            if y + bar_h > 7.2:
                break
            pct = rev / max_rev
            col = ACCENTS[i % len(ACCENTS)]
            bar_w = max(0.08, pct * bar_area_w)

            # label (left)
            _add_text(slide, mat[:32], 0.2, y, 3.8, bar_h,
                      size=9.5, color=C_SLATE, align=PP_ALIGN.RIGHT)
            # bar
            _add_rect(slide, 4.1, y + 0.04, bar_w, bar_h - 0.08, fill_hex=col)
            # value
            val_str = "₹" + (f"{rev/1e7:.2f} Cr" if rev >= 1e7 else f"{rev/1e5:.1f} L")
            _add_text(slide, val_str, 4.1 + bar_w + 0.08, y, 2.5, bar_h,
                      size=9, bold=True, color=col, align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════
    # SLIDE 4 — TOP CLIENTS / FACTORIES
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 7.5, fill_hex=C_BG)
    _add_rect(slide, 0, 0, 13.33, 0.75, fill_hex=C_GREEN)
    _add_text(slide, "Client Leaderboard", 0.3, 0.1, 9, 0.55,
              size=20, bold=True, color="#FFFFFF")
    _add_text(slide, "Top factories by total purchase value",
              0.3, 0.45, 9, 0.35, size=10, color="#D1FAE5")

    if not orders_df.empty:
        client_rev = (orders_df.groupby("customer_name")["total_amount"]
                      .sum().sort_values(ascending=False).head(10))
        client_qty = (items_df.groupby("customer_name")["quantity"]
                      .sum() if not items_df.empty else pd.Series())
        client_ord = orders_df.groupby("customer_name").size()
        total_r = client_rev.sum()

        # header
        _add_rect(slide, 0.25, 1.0, 12.85, 0.38, fill_hex=C_GREEN)
        for lbl, cx in zip(["#", "Factory / Client", "Revenue", "Share %", "Qty (MT)", "Orders"],
                            [0.25, 0.85, 5.2, 8.0, 9.8, 11.6]):
            _add_text(slide, lbl, cx + 0.05, 1.03, 1.8, 0.35,
                      size=10, bold=True, color="#FFFFFF")

        row_y = 1.38
        for rank, (client, rev) in enumerate(client_rev.items(), 1):
            bg = C_CARD if rank % 2 == 1 else "#F0FDF4"
            _add_rect(slide, 0.25, row_y, 12.85, 0.42, fill_hex=bg)
            pct = rev / total_r * 100
            qty = client_qty.get(client, 0)
            ord_ = client_ord.get(client, 0)
            rev_str = "₹" + (f"{rev/1e7:.2f} Cr" if rev >= 1e7 else f"{rev/1e5:.1f} L")
            row_data = [
                str(rank), client[:35], rev_str,
                f"{pct:.1f}%", f"{qty:,.1f}", str(ord_)
            ]
            col_xs2 = [0.25, 0.85, 5.2, 8.0, 9.8, 11.6]
            for val, cx in zip(row_data, col_xs2):
                bold_flag = val == str(rank)
                _add_text(slide, val, cx + 0.05, row_y + 0.05, 2.2, 0.35,
                          size=9.5, bold=(rank == 1 and val == rev_str),
                          color=(ACCENTS[(rank - 1) % len(ACCENTS)] if val == str(rank) else C_SLATE))
            row_y += 0.42

    # ══════════════════════════════════════════════════════════
    # SLIDE 5 — GEOGRAPHIC ANALYSIS
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 7.5, fill_hex=C_BG)
    _add_rect(slide, 0, 0, 13.33, 0.75, fill_hex=C_CYAN)
    _add_text(slide, "Geographic Analysis", 0.3, 0.1, 9, 0.55,
              size=20, bold=True, color=C_DARK)
    _add_text(slide, "State-wise revenue and order distribution",
              0.3, 0.45, 9, 0.35, size=10, color="#164E63")

    if not orders_df.empty:
        state_rev = (orders_df.groupby("consignee_state")["total_amount"]
                     .sum().sort_values(ascending=False))
        state_ord = orders_df.groupby("consignee_state").size()
        total_r2  = state_rev.sum()
        max_sr    = state_rev.max() if state_rev.max() > 0 else 1

        # header
        _add_rect(slide, 0.25, 1.0, 12.85, 0.38, fill_hex="#0891B2")
        for lbl, cx in zip(["State", "Revenue", "Share", "Orders", "Revenue Bar"],
                            [0.25, 3.8, 6.3, 8.2, 9.5]):
            _add_text(slide, lbl, cx + 0.05, 1.03, 2.5, 0.35,
                      size=10, bold=True, color="#FFFFFF")

        row_y = 1.38
        bar_max_w = 3.3
        for i, (state, rev) in enumerate(state_rev.items()):
            if row_y > 7.1:
                break
            bg = C_CARD if i % 2 == 0 else "#ECFEFF"
            _add_rect(slide, 0.25, row_y, 12.85, 0.37, fill_hex=bg)
            pct  = rev / total_r2 * 100
            ord_ = state_ord.get(state, 0)
            rev_str = "₹" + (f"{rev/1e7:.2f} Cr" if rev >= 1e7 else f"{rev/1e5:.1f} L")
            col = ACCENTS[i % len(ACCENTS)]
            bw  = max(0.05, (rev / max_sr) * bar_max_w)

            row_data = [state, rev_str, f"{pct:.1f}%", str(ord_)]
            col_xs3  = [0.25, 3.8, 6.3, 8.2]
            for val, cx in zip(row_data, col_xs3):
                _add_text(slide, val, cx + 0.05, row_y + 0.04, 2.5, 0.3,
                          size=9.5, color=C_SLATE)
            _add_rect(slide, 9.5, row_y + 0.08, bw, 0.22, fill_hex=col)
            row_y += 0.37

    # ══════════════════════════════════════════════════════════
    # SLIDE 6 — YEAR-WISE MATERIAL BREAKDOWN
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 7.5, fill_hex=C_BG)
    _add_rect(slide, 0, 0, 13.33, 0.75, fill_hex=C_PINK)
    _add_text(slide, "Material Mix by Financial Year", 0.3, 0.1, 9, 0.55,
              size=20, bold=True, color="#FFFFFF")
    _add_text(slide, "Top materials revenue contribution per FY",
              0.3, 0.45, 9, 0.35, size=10, color="#FCE7F3")

    if not items_df.empty:
        top5_mats = (items_df.groupby("material_name")["amount"]
                     .sum().sort_values(ascending=False).head(5).index.tolist())
        fy_list   = sorted(orders_df["financial_year"].dropna().unique())
        mat_fy    = (items_df[items_df["material_name"].isin(top5_mats)]
                     .groupby(["financial_year", "material_name"])["amount"]
                     .sum().unstack(fill_value=0))

        # header
        col_hs = ["Financial Year"] + top5_mats
        col_xs_m = [0.25] + [3.2 + i * 2.0 for i in range(len(top5_mats))]
        _add_rect(slide, 0.25, 1.0, 12.85, 0.38, fill_hex="#BE185D")
        for lbl, cx in zip(col_hs, col_xs_m):
            _add_text(slide, str(lbl)[:22], cx + 0.05, 1.03, 2.3, 0.35,
                      size=9, bold=True, color="#FFFFFF")

        row_y = 1.38
        for i, fy in enumerate(fy_list):
            bg = C_CARD if i % 2 == 0 else "#FDF2F8"
            _add_rect(slide, 0.25, row_y, 12.85, 0.42, fill_hex=bg)
            _add_text(slide, str(fy), col_xs_m[0] + 0.05, row_y + 0.06, 2.8, 0.32,
                      size=10, bold=True, color=C_DARK)
            for j, mat in enumerate(top5_mats):
                rev = mat_fy.loc[fy, mat] if fy in mat_fy.index and mat in mat_fy.columns else 0
                val_str = "₹" + (f"{rev/1e7:.2f} Cr" if rev >= 1e7 else f"{rev/1e5:.1f} L") if rev > 0 else "—"
                cx = col_xs_m[j + 1]
                _add_text(slide, val_str, cx + 0.05, row_y + 0.06, 2.0, 0.32,
                          size=9.5, color=(ACCENTS[j] if rev > 0 else C_MUTED))
            row_y += 0.42

    # ══════════════════════════════════════════════════════════
    # SLIDE 7 — THANK YOU / SUMMARY
    # ══════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank_layout)
    _add_rect(slide, 0, 0, 13.33, 7.5, fill_hex=C_DARK)
    _add_rect(slide, 0, 0, 0.55, 7.5, fill_hex=C_INDIGO)
    _add_rect(slide, 0.55, 3.5, 12.78, 0.05, fill_hex=C_INDIGO)

    _add_text(slide, "⚡ MS ENTERPRISES", 1.1, 1.6, 10, 0.7,
              size=13, bold=True, color=C_CYAN)
    _add_text(slide, "Data Summary", 1.1, 2.2, 10, 1.1,
              size=38, bold=True, color="#FFFFFF")
    _add_text(slide, "Confidential · Internal Use Only",
              1.1, 3.2, 10, 0.45, size=12, color="#94A3B8")

    summary_lines = [
        f"Total Revenue :  ₹{'%.2f Cr' % (total_rev/1e7) if total_rev >= 1e7 else '%.1f L' % (total_rev/1e5)}",
        f"Total Qty     :  {total_qty:,.1f} MT",
        f"Total Orders  :  {n_orders}",
        f"Factories     :  {n_cust}",
        f"States        :  {n_states}",
    ]
    sy = 4.1
    for i, line in enumerate(summary_lines):
        col = ACCENTS[i % len(ACCENTS)]
        _add_rect(slide, 1.1, sy, 0.04, 0.3, fill_hex=col)
        _add_text(slide, line, 1.25, sy - 0.02, 8, 0.38,
                  size=12, color="#CBD5E1")
        sy += 0.45

    _add_text(slide, f"Generated on {datetime.now().strftime('%d %b %Y at %H:%M')}",
              1.1, 6.8, 10, 0.4, size=9, color="#475569")

    # ── Save to bytes ──────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()

DB_PATH = os.path.join(os.path.dirname(__file__), "ms_enterprises.db")

st.set_page_config(
    page_title="MS Enterprises · PO Intelligence",
    page_icon="⚡",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── THEME STATE (must be before CSS injection) ──────────────
dark_mode = False  # white theme always

# ═══════════════════════════════════════════════════════════
# PREMIUM CSS — DARK GLASSMORPHISM INDUSTRIAL THEME
# ═══════════════════════════════════════════════════════════
PREMIUM_CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap');

:root {
  --bg:  #070C18;
  --bg2: #0D1426;
  --glass: rgba(13,22,45,0.75);
  --p:  #6366F1;
  --c:  #22D3EE;
  --o:  #F59E0B;
  --g:  #10B981;
  --r:  #EF4444;
  --pk: #EC4899;
  --t1: #F1F5F9;
  --t2: #94A3B8;
  --t3: #475569;
  --br: rgba(255,255,255,0.07);
  --r1: 18px;
  --r2: 12px;
  --tr: all 0.3s cubic-bezier(0.4,0,0.2,1);
  --sh: 0 8px 32px rgba(0,0,0,0.45);
  --sh2: 0 20px 60px rgba(0,0,0,0.6);
}

html,body,[class*="css"],.stApp {
  font-family:'Inter',-apple-system,BlinkMacSystemFont,sans-serif !important;
  background: var(--bg) !important;
  color: var(--t1) !important;
}
.stApp {
  background: radial-gradient(ellipse at 20% 0%, rgba(99,102,241,0.12) 0%, transparent 50%),
              radial-gradient(ellipse at 80% 100%, rgba(34,211,238,0.08) 0%, transparent 50%),
              var(--bg) !important;
  min-height: 100vh;
}

#MainMenu,footer,[data-testid="stToolbar"],
[data-testid="stDecoration"],[data-testid="stStatusWidget"] { display:none !important; }
header[data-testid="stHeader"] { background:transparent !important; height:0 !important; }

.main .block-container {
  padding: 1.5rem 2rem 3rem !important;
  max-width: 100% !important;
  background: transparent !important;
}

/* ── Sidebar ── */
[data-testid="stSidebar"] {
  background: linear-gradient(180deg,#0A0F22 0%,#070C18 100%) !important;
  border-right: 1px solid var(--br) !important;
  width: 252px !important;
}
[data-testid="stSidebar"] > div { padding-top:0 !important; }
[data-testid="stSidebarContent"] { padding:0 !important; }

/* ── Nav radio ── */
.stRadio > div { gap:3px !important; }
.stRadio label {
  background:transparent !important;
  border-radius:var(--r2) !important;
  padding:9px 16px !important;
  cursor:pointer !important;
  transition:var(--tr) !important;
  color:var(--t2) !important;
  font-size:13.5px !important;
  font-weight:500 !important;
  border:1px solid transparent !important;
  display:flex !important; align-items:center !important;
}
.stRadio label:hover {
  background:rgba(99,102,241,0.12) !important;
  color:var(--t1) !important;
  border-color:rgba(99,102,241,0.25) !important;
}
.stRadio input[type="radio"] { display:none !important; }
.stRadio [data-testid="stRadioLabel"] > div:first-child { display:none !important; }

/* ── Multiselect ── */
.stMultiSelect [data-baseweb="tag"] {
  background:rgba(99,102,241,0.2) !important;
  border:1px solid rgba(99,102,241,0.4) !important;
  border-radius:999px !important;
  color:#A5B4FC !important;
  font-size:12px !important;
}
.stMultiSelect [data-baseweb="select"],
.stSelectbox [data-baseweb="select"] {
  background:rgba(255,255,255,0.05) !important;
  border:1px solid var(--br) !important;
  border-radius:var(--r2) !important;
  color:var(--t1) !important;
}

/* ── Buttons ── */
.stButton > button {
  background:linear-gradient(135deg,var(--p),#4338CA) !important;
  color:white !important; border:none !important;
  border-radius:var(--r2) !important;
  padding:8px 20px !important; font-weight:600 !important;
  box-shadow:0 4px 15px rgba(99,102,241,0.35) !important;
  transition:var(--tr) !important;
}
.stButton > button:hover {
  transform:translateY(-2px) !important;
  box-shadow:0 8px 25px rgba(99,102,241,0.5) !important;
}
.stDownloadButton > button {
  background:rgba(34,211,238,0.12) !important;
  color:var(--c) !important;
  border:1px solid rgba(34,211,238,0.3) !important;
  border-radius:var(--r2) !important; font-weight:600 !important;
}

/* ── Inputs ── */
.stTextInput input,.stTextArea textarea {
  background:rgba(255,255,255,0.06) !important;
  border:1px solid var(--br) !important;
  border-radius:var(--r2) !important;
  color:var(--t1) !important; font-size:14px !important;
}
.stTextInput input:focus,.stTextArea textarea:focus {
  border-color:var(--p) !important;
  box-shadow:0 0 0 3px rgba(99,102,241,0.2) !important;
}

/* ── Dataframe ── */
[data-testid="stDataFrame"] {
  border-radius:var(--r1) !important;
  overflow:hidden !important;
  border:1px solid var(--br) !important;
  background:var(--glass) !important;
}
[data-testid="stDataFrame"] th {
  background:rgba(99,102,241,0.15) !important;
  color:var(--t1) !important; font-weight:600 !important;
}
[data-testid="stDataFrame"] td { color:var(--t2) !important; }

/* ── Plotly ── */
[data-testid="stPlotlyChart"] {
  background:var(--glass) !important;
  backdrop-filter:blur(16px) !important;
  border:1px solid var(--br) !important;
  border-radius:var(--r1) !important;
  padding:4px !important;
  box-shadow:var(--sh) !important;
  transition:var(--tr) !important;
}
[data-testid="stPlotlyChart"]:hover {
  border-color:rgba(99,102,241,0.3) !important;
  box-shadow:var(--sh2) !important;
}

/* ── Scrollbar ── */
::-webkit-scrollbar { width:5px; height:5px; }
::-webkit-scrollbar-track { background:transparent; }
::-webkit-scrollbar-thumb { background:rgba(255,255,255,0.1); border-radius:99px; }
::-webkit-scrollbar-thumb:hover { background:rgba(99,102,241,0.5); }

hr { border-color:var(--br) !important; }

/* ══════ KPI CARDS (external CSS for in-page use) ══════ */
.kpi-grid {
  display:grid; grid-template-columns:repeat(5,1fr); gap:14px; margin-bottom:24px;
}

/* ══════ SECTION HEADER ══════ */
.sec-header {
  display:flex; align-items:center; gap:10px;
  font-size:17px; font-weight:700; color:var(--t1);
  margin:28px 0 16px; letter-spacing:-0.01em;
}
.sec-header .dot {
  width:7px; height:7px; border-radius:50%;
  background:var(--p); box-shadow:0 0 8px var(--p);
}
.sec-sub { font-size:12px; color:var(--t3); margin-top:-10px; margin-bottom:16px; }

/* ══════ PAGE TITLE ══════ */
.page-title {
  font-size:28px; font-weight:800; color:var(--t1);
  letter-spacing:-0.03em; margin-bottom:4px;
}
.page-title span {
  background:linear-gradient(135deg,var(--p),var(--c));
  -webkit-background-clip:text; -webkit-text-fill-color:transparent;
}
.page-sub { font-size:13px; color:var(--t3); margin-bottom:20px; }

/* ══════ INSIGHT CARDS ══════ */
.insight-grid {
  display:grid; grid-template-columns:repeat(3,1fr); gap:12px; margin-bottom:20px;
}
.insight-card {
  background:rgba(255,255,255,0.03);
  border:1px solid var(--br);
  border-left:3px solid var(--ic-color);
  border-radius:var(--r2); padding:14px 16px;
  backdrop-filter:blur(12px); transition:var(--tr);
}
.insight-card:hover { background:rgba(255,255,255,0.06); transform:translateX(3px); }
.insight-card .ic-icon { font-size:20px; margin-bottom:8px; }
.insight-card .ic-title { font-size:11px; font-weight:600; color:var(--t3); text-transform:uppercase; letter-spacing:.06em; margin-bottom:4px; }
.insight-card .ic-value { font-size:14px; font-weight:700; color:var(--t1); }
.insight-card .ic-sub { font-size:12px; color:var(--t2); margin-top:2px; }

/* ══════ AI PANEL ══════ */
.ai-panel {
  background:linear-gradient(135deg,rgba(99,102,241,0.1),rgba(34,211,238,0.06));
  border:1px solid rgba(99,102,241,0.25);
  border-radius:var(--r1); padding:20px 24px;
  position:relative; overflow:hidden;
  backdrop-filter:blur(16px);
}
.ai-panel::before {
  content:'✦'; position:absolute; top:14px; right:18px;
  font-size:22px; color:rgba(99,102,241,0.35);
  animation:pulse 2s ease-in-out infinite;
}
@keyframes pulse { 0%,100%{opacity:.3;transform:scale(1)} 50%{opacity:.9;transform:scale(1.2)} }
.ai-panel h4 { font-size:13px; font-weight:700; color:#A5B4FC; margin:0 0 12px; }
.ai-bullet {
  display:flex; align-items:flex-start; gap:10px;
  padding:8px 0; border-bottom:1px solid rgba(255,255,255,0.05);
  font-size:13px; color:var(--t2);
}
.ai-bullet:last-child { border-bottom:none; }
.ai-bullet .bul { color:var(--p); font-weight:700; flex-shrink:0; }
.ai-bullet strong { color:var(--t1); }

/* ══════ LEADERBOARD ══════ */
.lb-item {
  display:flex; align-items:center; gap:12px;
  padding:12px 16px;
  background:rgba(255,255,255,0.03);
  border:1px solid var(--br); border-radius:var(--r2);
  margin-bottom:7px; transition:var(--tr);
  backdrop-filter:blur(12px);
}
.lb-item:hover { background:rgba(99,102,241,0.08); border-color:rgba(99,102,241,0.25); }
.lb-rank {
  width:28px; height:28px; border-radius:8px;
  background:rgba(99,102,241,0.15); color:#A5B4FC;
  font-size:12px; font-weight:700;
  display:flex; align-items:center; justify-content:center; flex-shrink:0;
}
.lb-rank.gold   { background:rgba(245,158,11,0.18); color:#F59E0B; }
.lb-rank.silver { background:rgba(148,163,184,0.15); color:#94A3B8; }
.lb-rank.bronze { background:rgba(180,120,60,0.15); color:#B47C3C; }
.lb-name { flex:1; font-size:13px; font-weight:600; color:var(--t1); }
.lb-sub  { font-size:11px; color:var(--t3); margin-top:2px; }
.lb-rev  { font-size:14px; font-weight:700; color:var(--t1); text-align:right; }
.lb-bar-wrap { width:80px; height:3px; background:rgba(255,255,255,0.07); border-radius:99px; margin-top:4px; }
.lb-bar { height:3px; border-radius:99px; background:linear-gradient(90deg,var(--p),var(--c)); }

/* ══════ SIDEBAR LOGO ══════ */
.sb-logo {
  padding:22px 20px 18px;
  border-bottom:1px solid var(--br);
  margin-bottom:6px;
  background:linear-gradient(135deg,rgba(99,102,241,0.08),transparent);
}
.sb-logo-title {
  font-size:17px; font-weight:800; letter-spacing:-0.02em;
  background:linear-gradient(135deg,#A5B4FC,#22D3EE);
  -webkit-background-clip:text; -webkit-text-fill-color:transparent;
}
.sb-logo-sub { font-size:11px; color:var(--t3); margin-top:2px; }
.sb-logo-badge {
  display:inline-block; padding:2px 8px;
  background:rgba(99,102,241,0.2);
  border:1px solid rgba(99,102,241,0.35);
  border-radius:999px; font-size:10px; color:#A5B4FC;
  font-weight:600; margin-top:8px; letter-spacing:.05em;
}
.sb-nav-label {
  font-size:10px; color:var(--t3); font-weight:700;
  letter-spacing:.1em; text-transform:uppercase;
  padding:14px 20px 8px;
}
.sb-footer {
  padding:14px 20px;
  border-top:1px solid var(--br);
  font-size:11px; color:var(--t3);
}
.sb-footer strong { color:var(--t2); }

/* ══════ BADGES ══════ */
.badge { display:inline-flex; align-items:center; gap:5px; padding:3px 10px; border-radius:999px; font-size:11px; font-weight:600; }
.badge-green  { background:rgba(16,185,129,0.15); color:#10B981; }
.badge-red    { background:rgba(239,68,68,0.15); color:#EF4444; }
.badge-yellow { background:rgba(245,158,11,0.15); color:#F59E0B; }
.badge-blue   { background:rgba(99,102,241,0.15); color:#A5B4FC; }

/* ══════ GLASS CARD ══════ */
.glass-card {
  background:var(--glass);
  backdrop-filter:blur(20px);
  border:1px solid var(--br);
  border-radius:var(--r1);
  padding:20px;
  box-shadow:var(--sh);
}

/* ══════ FILTER BAR ══════ */
.filter-bar {
  background:rgba(13,22,45,0.85);
  backdrop-filter:blur(20px);
  border:1px solid var(--br);
  border-radius:var(--r1);
  padding:12px 20px;
  margin-bottom:20px;
  display:flex; gap:12px; align-items:center; flex-wrap:wrap;
}
.filter-label {
  font-size:11px; font-weight:700; text-transform:uppercase;
  letter-spacing:.08em; color:var(--t3);
}

/* ══════ METRIC MINI ══════ */
.metric-mini {
  background:var(--glass); border:1px solid var(--br);
  border-radius:var(--r2); padding:16px;
  backdrop-filter:blur(16px);
}
.metric-mini .mm-val { font-size:20px; font-weight:700; color:var(--t1); }
.metric-mini .mm-lbl { font-size:11px; color:var(--t3); font-weight:600; text-transform:uppercase; letter-spacing:.06em; margin-top:4px; }

/* ══════ ANIMATIONS ══════ */
@keyframes fadeUp { from{opacity:0;transform:translateY(14px)} to{opacity:1;transform:translateY(0)} }
@keyframes shimmer { 0%{background-position:-200% 0} 100%{background-position:200% 0} }
.fade-up { animation:fadeUp .45s ease both; }
.fade-up:nth-child(1){animation-delay:.05s}
.fade-up:nth-child(2){animation-delay:.10s}
.fade-up:nth-child(3){animation-delay:.15s}
.fade-up:nth-child(4){animation-delay:.20s}
.fade-up:nth-child(5){animation-delay:.25s}
.skeleton {
  background:linear-gradient(90deg,rgba(255,255,255,0.05) 25%,rgba(255,255,255,0.1) 50%,rgba(255,255,255,0.05) 75%);
  background-size:200% 100%; animation:shimmer 1.5s infinite; border-radius:8px;
}
</style>
"""
LIGHT_OVERRIDES = """<style>
:root {
  --bg:#F8FAFC; --bg2:#FFFFFF; --glass:rgba(255,255,255,0.92);
  --t1:#0F172A; --t2:#334155; --t3:#64748B;
  --br:rgba(15,23,42,0.08);
  --sh:0 2px 20px rgba(15,23,42,0.08);
  --sh2:0 8px 40px rgba(15,23,42,0.14);
}
html,body,[class*="css"],.stApp{background:#F8FAFC !important;color:#0F172A !important;}
.stApp{
  background:
    radial-gradient(ellipse at 15% 0%,rgba(99,102,241,0.06) 0%,transparent 45%),
    radial-gradient(ellipse at 85% 100%,rgba(34,211,238,0.04) 0%,transparent 45%),
    #F8FAFC !important;
}
[data-testid="stSidebar"]{
  background:linear-gradient(180deg,#FFFFFF 0%,#F1F5F9 100%) !important;
  border-right:1px solid rgba(15,23,42,0.08) !important;
}
[data-testid="stSidebar"] > div > div,[data-testid="stSidebarContent"]{background:transparent !important;}
.sb-logo{background:linear-gradient(135deg,rgba(99,102,241,0.06),transparent) !important;}
.sb-logo-title{background:linear-gradient(135deg,#6366F1,#22D3EE);-webkit-background-clip:text;-webkit-text-fill-color:transparent;}
.sb-logo-sub,.sb-footer{color:#64748B !important;}
.sb-nav-label{color:#94A3B8 !important;}
.sb-logo-badge{background:rgba(99,102,241,0.1) !important;border-color:rgba(99,102,241,0.2) !important;color:#6366F1 !important;}
.stRadio label{color:#475569 !important;border-color:transparent !important;}
.stRadio label:hover{background:rgba(99,102,241,0.07) !important;color:#0F172A !important;border-color:rgba(99,102,241,0.2) !important;}
[data-testid="stDataFrame"]{background:white !important;border-color:rgba(0,0,0,0.08) !important;}
[data-testid="stDataFrame"] th{background:rgba(99,102,241,0.08) !important;color:#0F172A !important;}
[data-testid="stDataFrame"] td{color:#334155 !important;}
[data-testid="stPlotlyChart"]{background:white !important;border-color:rgba(15,23,42,0.08) !important;box-shadow:0 2px 20px rgba(15,23,42,0.08) !important;}
[data-testid="stPlotlyChart"]:hover{border-color:rgba(99,102,241,0.3) !important;box-shadow:0 8px 32px rgba(99,102,241,0.12) !important;}
.lb-item{background:rgba(255,255,255,0.9) !important;border-color:rgba(15,23,42,0.07) !important;}
.lb-item:hover{background:rgba(99,102,241,0.05) !important;}
.lb-name{color:#0F172A !important;}.lb-sub{color:#64748B !important;}.lb-rev{color:#0F172A !important;}
.lb-rank{background:rgba(99,102,241,0.1) !important;color:#6366F1 !important;}
.lb-rank.gold{background:rgba(245,158,11,0.12) !important;color:#D97706 !important;}
.lb-rank.silver{background:rgba(100,116,139,0.12) !important;color:#475569 !important;}
.lb-rank.bronze{background:rgba(120,80,40,0.12) !important;color:#92400E !important;}
.insight-card{background:rgba(255,255,255,0.85) !important;border-color:rgba(15,23,42,0.07) !important;}
.insight-card:hover{background:rgba(255,255,255,0.97) !important;}
.ic-value{color:#0F172A !important;}.ic-sub{color:#475569 !important;}
.ai-panel{background:linear-gradient(135deg,rgba(99,102,241,0.05),rgba(34,211,238,0.03)) !important;border-color:rgba(99,102,241,0.18) !important;}
.ai-panel h4{color:#6366F1 !important;}.ai-bullet{color:#334155 !important;}
.ai-bullet strong{color:#0F172A !important;}
.glass-card{background:rgba(255,255,255,0.92) !important;border-color:rgba(15,23,42,0.08) !important;}
.filter-bar{background:rgba(255,255,255,0.95) !important;border-color:rgba(15,23,42,0.08) !important;}
.metric-mini{background:white !important;border-color:rgba(15,23,42,0.07) !important;}
.metric-mini .mm-val{color:#0F172A !important;}.metric-mini .mm-lbl{color:#64748B !important;}
.page-title{color:#0F172A !important;}.sec-header{color:#0F172A !important;}
.sec-header .dot{background:#6366F1 !important;box-shadow:0 0 8px rgba(99,102,241,0.4) !important;}
.page-sub{color:#64748B !important;}.sec-sub{color:#94A3B8 !important;}
.badge-green{background:rgba(16,185,129,0.1) !important;color:#059669 !important;}
.badge-red{background:rgba(239,68,68,0.1) !important;color:#DC2626 !important;}
.badge-blue{background:rgba(99,102,241,0.1) !important;color:#6366F1 !important;}
hr{border-color:rgba(15,23,42,0.08) !important;}
::-webkit-scrollbar-thumb{background:rgba(15,23,42,0.12) !important;}
.stTextInput input,.stTextArea textarea{
  background:white !important;border-color:rgba(15,23,42,0.1) !important;
  color:#0F172A !important;box-shadow:none !important;
}
.stTextInput input:focus,.stTextArea textarea:focus{border-color:#6366F1 !important;}
.stMultiSelect [data-baseweb="select"],.stSelectbox [data-baseweb="select"]{
  background:white !important;border-color:rgba(15,23,42,0.1) !important;color:#0F172A !important;
}
.stMultiSelect [data-baseweb="tag"]{
  background:rgba(99,102,241,0.1) !important;border-color:rgba(99,102,241,0.2) !important;color:#6366F1 !important;
}
[data-testid="stMetricValue"]{color:#0F172A !important;}
[data-testid="stMetricLabel"]{color:#64748B !important;}
.stButton > button{box-shadow:0 4px 12px rgba(99,102,241,0.25) !important;}
</style>"""

st.markdown(PREMIUM_CSS, unsafe_allow_html=True)
if not dark_mode:
    st.markdown(LIGHT_OVERRIDES, unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════
ACCENT = ["#6366F1","#22D3EE","#F59E0B","#10B981","#EC4899","#F97316","#A78BFA","#EF4444"]

if dark_mode:
    _txt   = "#94A3B8"
    _tick  = "#475569"
    _grid  = "rgba(255,255,255,0.05)"
    _zero  = "rgba(255,255,255,0.04)"
    _hbg   = "#0D1426"
    _hbdr  = "rgba(255,255,255,0.1)"
    _hfnt  = "#F1F5F9"
    _tfnt  = "#F1F5F9"
    _leg   = "rgba(13,22,45,0.85)"
    _legbc = "rgba(255,255,255,0.07)"
else:
    _txt   = "#334155"
    _tick  = "#64748B"
    _grid  = "rgba(0,0,0,0.05)"
    _zero  = "rgba(0,0,0,0.06)"
    _hbg   = "#FFFFFF"
    _hbdr  = "rgba(0,0,0,0.1)"
    _hfnt  = "#0F172A"
    _tfnt  = "#0F172A"
    _leg   = "rgba(255,255,255,0.95)"
    _legbc = "rgba(0,0,0,0.08)"

CHART_LAYOUT = dict(
    paper_bgcolor="rgba(0,0,0,0)", plot_bgcolor="rgba(0,0,0,0)",
    font=dict(family="Inter, sans-serif", color=_txt, size=12),
    margin=dict(l=10, r=10, t=44, b=10),
    showlegend=True,
    legend=dict(bgcolor=_leg, bordercolor=_legbc,
                borderwidth=1, font=dict(color=_txt, size=11)),
    xaxis=dict(gridcolor=_grid, zerolinecolor=_zero,
               tickfont=dict(color=_tick, size=11), showline=False),
    yaxis=dict(gridcolor=_grid, zerolinecolor=_zero,
               tickfont=dict(color=_tick, size=11), showline=False),
    hoverlabel=dict(bgcolor=_hbg, bordercolor=_hbdr,
                    font=dict(color=_hfnt, size=12)),
    title_font=dict(color=_tfnt, size=15, family="Inter, sans-serif"),
    title_x=0.02,
)

# ═══════════════════════════════════════════════════════════
# SESSION STATE
# ═══════════════════════════════════════════════════════════
if "page" not in st.session_state:
    st.session_state.page = "📊  Overview"
if "fy_filter" not in st.session_state:
    st.session_state.fy_filter = []
if "state_filter" not in st.session_state:
    st.session_state.state_filter = []

# ═══════════════════════════════════════════════════════════
# DATA LOADING
# ═══════════════════════════════════════════════════════════
@st.cache_data(ttl=300)
def load_data():
    conn = sqlite3.connect(DB_PATH)
    orders = pd.read_sql_query("""
        SELECT o.*, o.consignee_name as customer_name,
               o.supplier_name as vendor_name,
               o.consignee_state as delivery_state
        FROM orders o
    """, conn)
    items = pd.read_sql_query("""
        SELECT oi.*, m.name_standardized as material_name,
               m.category as material_category,
               o.financial_year, o.document_date, o.consignee_state,
               o.consignee_name as customer_name,
               o.supplier_name as vendor_name,
               o.month, o.year, o.po_number
        FROM order_items oi
        JOIN orders o ON oi.order_id = o.id
        LEFT JOIN materials m ON oi.material_id = m.id
    """, conn)
    processing = pd.read_sql_query("SELECT * FROM processing_log", conn)
    conn.close()

    orders["document_date"] = pd.to_datetime(orders["document_date"], errors="coerce")
    items["document_date"]  = pd.to_datetime(items["document_date"], errors="coerce")

    # ── Unit normalisation: convert KG → MT; mark non-weight units ──
    # Items in KG are divided by 1000 to get MT; rate adjusted ×1000
    if "unit" in items.columns:
        kg_mask  = items["unit"].str.upper().eq("KG")
        non_wt   = items["unit"].str.upper().isin(["PCS","EA","NOS","DRUM","SET","LOT"])
        # KG items: qty_mt = qty/1000, rate stays per-KG so adjust to per-MT
        items.loc[kg_mask, "quantity"] = items.loc[kg_mask, "quantity"] / 1000
        items.loc[kg_mask, "rate"]     = items.loc[kg_mask, "rate"]     * 1000
        # Non-weight items: zero out quantity so they don't inflate MT totals
        items.loc[non_wt, "quantity"]  = 0

    return orders, items, processing


@st.cache_data(ttl=300)
def compute_insights(orders_df, items_df):
    insights = []
    if orders_df.empty:
        return insights

    # YoY revenue growth
    fy_rev = orders_df.groupby("financial_year")["total_amount"].sum().sort_index()
    if len(fy_rev) >= 2:
        latest_fy = fy_rev.index[-1]
        prev_fy   = fy_rev.index[-2]
        growth = (fy_rev[latest_fy] - fy_rev[prev_fy]) / fy_rev[prev_fy] * 100
        arrow = "↑" if growth > 0 else "↓"
        insights.append({
            "icon":"📈","color":"#10B981",
            "text": f"Revenue {arrow} <strong>{abs(growth):.1f}%</strong> in {latest_fy} vs {prev_fy} (₹{fy_rev[latest_fy]/1e7:.2f} Cr vs ₹{fy_rev[prev_fy]/1e7:.2f} Cr)"
        })

    # Top customer
    top_cust = orders_df.groupby("customer_name")["total_amount"].sum()
    if not top_cust.empty:
        tc = top_cust.idxmax()
        tc_pct = top_cust.max() / orders_df["total_amount"].sum() * 100
        tc_rev = top_cust.max()
        insights.append({
            "icon":"🏆","color":"#F59E0B",
            "text": f"<strong>{tc}</strong> is top customer — contributes <strong>{tc_pct:.1f}%</strong> of total revenue (₹{tc_rev/1e5:.1f} L)"
        })

    # Top state
    top_state = orders_df.groupby("consignee_state")["total_amount"].sum()
    if not top_state.empty:
        ts = top_state.idxmax()
        ts_pct = top_state.max() / orders_df["total_amount"].sum() * 100
        insights.append({
            "icon":"🗺️","color":"#06B6D4",
            "text": f"<strong>{ts}</strong> is highest revenue state — <strong>{ts_pct:.1f}%</strong> share across {orders_df.groupby('consignee_state').size()[ts]} orders"
        })

    # Top material
    if not items_df.empty:
        top_mat = items_df.groupby("material_name")["amount"].sum()
        if not top_mat.empty:
            tm = top_mat.idxmax()
            tm_pct = top_mat.max() / items_df["amount"].sum() * 100
            insights.append({
                "icon":"🧱","color":"#F97316",
                "text": f"<strong>{tm}</strong> is best-selling material — <strong>{tm_pct:.1f}%</strong> of item-level revenue"
            })

    # Best month
    month_rev = orders_df.dropna(subset=["document_date"])
    if not month_rev.empty:
        month_rev = month_rev.copy()
        month_rev["mn"] = month_rev["document_date"].dt.strftime("%B")
        month_rev["m"]  = month_rev["document_date"].dt.month
        bm = month_rev.groupby(["m","mn"])["total_amount"].sum().reset_index()
        best = bm.loc[bm["total_amount"].idxmax()]
        insights.append({
            "icon":"📅","color":"#EC4899",
            "text": f"<strong>{best['mn']}</strong> is historically highest revenue month (₹{best['total_amount']/1e5:.1f} L avg across all years)"
        })

    # Avg order size
    avg = orders_df["total_amount"].mean()
    max_o = orders_df["total_amount"].max()
    insights.append({
        "icon":"💡","color":"#7C3AED",
        "text": f"Average order value <strong>₹{avg/1e5:.2f} L</strong> — largest single order was <strong>₹{max_o/1e5:.1f} L</strong>"
    })

    return insights


# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════
def fmt(amount):
    if not amount or np.isnan(amount): return "₹0"
    if amount >= 1e7:  return f"₹{amount/1e7:.2f} Cr"
    if amount >= 1e5:  return f"₹{amount/1e5:.2f} L"
    if amount >= 1e3:  return f"₹{amount/1e3:.1f}K"
    return f"₹{amount:,.0f}"


def spark_svg(values, color="#7C3AED", w=80, h=28):
    if not values: return ""
    vals = [v for v in values if v is not None and not np.isnan(float(v))]
    if len(vals) < 2: return ""
    vmin, vmax = min(vals), max(vals)
    rng = vmax - vmin or 1
    pts = []
    for i, v in enumerate(vals):
        x = i * (w / (len(vals)-1))
        y = h - ((v - vmin) / rng) * (h - 4) - 2
        pts.append(f"{x:.1f},{y:.1f}")
    path = "M " + " L ".join(pts)
    fill = f"{path} L {w},{h} L 0,{h} Z"
    uid = color.replace("#","")
    return f"""<svg width="{w}" height="{h}" viewBox="0 0 {w} {h}" style="display:block">
  <defs>
    <linearGradient id="sg{uid}" x1="0" y1="0" x2="0" y2="1">
      <stop offset="0%" stop-color="{color}" stop-opacity="0.35"/>
      <stop offset="100%" stop-color="{color}" stop-opacity="0"/>
    </linearGradient>
  </defs>
  <path d="{fill}" fill="url(#sg{uid})"/>
  <path d="{path}" stroke="{color}" stroke-width="2" fill="none" stroke-linecap="round" stroke-linejoin="round"/>
</svg>"""


def kpi_card_html(icon, label, value_num, display_val, trend_pct, spark_data,
                  accent, accent2, glow, fmt_type="currency"):
    """Returns self-contained HTML for one KPI card (used inside components.html)."""
    trend_cls   = "up" if trend_pct > 0 else ("down" if trend_pct < 0 else "flat")
    trend_arrow = "↑" if trend_pct > 0 else ("↓" if trend_pct < 0 else "→")
    trend_html  = (f'<span class="kpi-trend {trend_cls}">'
                   f'{trend_arrow} {abs(trend_pct):.1f}%</span>') if trend_pct != 0 else "<span></span>"
    spk = spark_svg(spark_data, accent)
    return f"""<div class="kpi-card" style="--ac:{accent};--ac2:{accent2};--glow:{glow}">
  <div class="kpi-top">
    <div class="kpi-icon">{icon}</div>
    <div class="kpi-label">{label}</div>
  </div>
  <div class="kpi-value" data-target="{value_num}" data-fmt="{fmt_type}">{display_val}</div>
  <div class="kpi-footer">{trend_html}<div class="kpi-spark">{spk}</div></div>
</div>"""


def render_kpi_hero(cards_data):
    """Render the full KPI hero row in an iframe (avoids Streamlit markdown parser)."""
    cards_html = "\n".join([kpi_card_html(*c) for c in cards_data])
    _card_bg   = "rgba(13,22,45,0.82)" if dark_mode else "rgba(255,255,255,0.95)"
    _card_bdr  = "rgba(255,255,255,0.08)" if dark_mode else "rgba(0,0,0,0.08)"
    _card_sh   = "0 8px 32px rgba(0,0,0,0.45)" if dark_mode else "0 4px 20px rgba(0,0,0,0.10)"
    _card_val  = "#F1F5F9" if dark_mode else "#0F172A"
    _card_lbl  = "#475569" if dark_mode else "#64748B"
    _body_bg   = "transparent" if dark_mode else "#F8FAFC"
    full_html = f"""<!DOCTYPE html>
<html>
<head>
<style>
  @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700;800&display=swap');
  * {{ box-sizing:border-box; margin:0; padding:0; }}
  body {{ background:{_body_bg}; font-family:'Inter',sans-serif; }}
  .kpi-grid {{
    display:grid; grid-template-columns:repeat(5,1fr); gap:14px; padding:2px;
  }}
  .kpi-card {{
    position:relative;
    background:{_card_bg};
    backdrop-filter:blur(20px); -webkit-backdrop-filter:blur(20px);
    border:1px solid {_card_bdr}; border-radius:18px;
    padding:20px 18px 14px; overflow:hidden;
    box-shadow:{_card_sh};
    transition:transform .3s ease,box-shadow .3s ease,border-color .3s ease;
    cursor:default;
  }}
  .kpi-card::before {{
    content:''; position:absolute; top:0; left:0; right:0; height:2px;
    background:linear-gradient(90deg,var(--ac),var(--ac2)); opacity:.95;
  }}
  .kpi-card::after {{
    content:''; position:absolute; top:-45px; right:-45px;
    width:120px; height:120px;
    background:radial-gradient(circle,var(--ac) 0%,transparent 70%);
    opacity:.08; border-radius:50%; transition:opacity .3s,transform .3s;
  }}
  .kpi-card:hover {{
    transform:translateY(-5px); border-color:var(--ac);
    box-shadow:0 20px 50px rgba(0,0,0,.6),0 0 0 1px var(--ac);
  }}
  .kpi-card:hover::after {{ opacity:.18; transform:scale(1.25); }}
  .kpi-top {{ display:flex; align-items:center; gap:10px; margin-bottom:14px; }}
  .kpi-icon {{
    width:36px; height:36px; border-radius:10px; flex-shrink:0;
    background:linear-gradient(135deg,var(--ac),var(--ac2));
    display:flex; align-items:center; justify-content:center;
    font-size:17px; box-shadow:0 4px 12px var(--glow);
  }}
  .kpi-label {{
    font-size:10px; font-weight:700; letter-spacing:.08em;
    text-transform:uppercase; color:{_card_lbl};
  }}
  .kpi-value {{
    font-size:24px; font-weight:800; color:{_card_val};
    letter-spacing:-.02em; line-height:1; margin-bottom:12px;
    font-variant-numeric:tabular-nums;
  }}
  .kpi-footer {{ display:flex; align-items:center; justify-content:space-between; }}
  .kpi-trend {{
    display:inline-flex; align-items:center; gap:3px;
    font-size:11px; font-weight:700; padding:2px 8px; border-radius:999px;
  }}
  .kpi-trend.up   {{ background:rgba(16,185,129,.15); color:#10B981; }}
  .kpi-trend.down {{ background:rgba(239,68,68,.15);  color:#EF4444; }}
  .kpi-trend.flat {{ background:rgba(71,85,105,.2);   color:#94A3B8; }}
  .kpi-spark {{ opacity:.75; display:flex; align-items:flex-end; }}
  @keyframes fadeUp {{
    from {{ opacity:0; transform:translateY(14px); }}
    to   {{ opacity:1; transform:translateY(0); }}
  }}
  .kpi-card {{ animation:fadeUp .45s ease both; }}
  .kpi-card:nth-child(1) {{ animation-delay:.05s; }}
  .kpi-card:nth-child(2) {{ animation-delay:.10s; }}
  .kpi-card:nth-child(3) {{ animation-delay:.15s; }}
  .kpi-card:nth-child(4) {{ animation-delay:.20s; }}
  .kpi-card:nth-child(5) {{ animation-delay:.25s; }}
</style>
</head>
<body>
<div class="kpi-grid">
{cards_html}
</div>
<script>
(function(){{
  document.querySelectorAll('.kpi-value[data-target]').forEach(function(el){{
    var raw = parseInt(el.getAttribute('data-target'));
    var fmt = el.getAttribute('data-fmt') || 'currency';
    if(!raw || isNaN(raw)) return;
    var start = performance.now(), dur = 1500;
    function tick(now){{
      var p = Math.min((now-start)/dur, 1);
      var e = 1 - Math.pow(1-p, 3);
      var v = e * raw;
      if(fmt === 'currency') {{
        if(raw >= 10000000)     el.textContent = '₹' + (v/10000000).toFixed(2) + ' Cr';
        else if(raw >= 100000)  el.textContent = '₹' + (v/100000).toFixed(1) + ' L';
        else if(raw >= 1000)    el.textContent = '₹' + Math.round(v).toLocaleString('en-IN');
        else                    el.textContent = '₹' + Math.round(v);
      }} else if(fmt === 'qty') {{
        if(raw >= 100000)       el.textContent = (v/100000).toFixed(1) + ' L MT';
        else if(raw >= 1000)    el.textContent = (v/1000).toFixed(1) + 'K MT';
        else                    el.textContent = v.toFixed(1) + ' MT';
      }} else {{
        el.textContent = Math.round(v).toLocaleString('en-IN');
      }}
      if(p < 1) requestAnimationFrame(tick);
    }}
    requestAnimationFrame(tick);
  }});
}})();
</script>
</body>
</html>"""
    components.html(full_html, height=170, scrolling=False)


def section_header(title, subtitle=""):
    sub = f'<div class="page-sub">{subtitle}</div>' if subtitle else ""
    return f'<div class="sec-header"><div class="dot"></div>{title}</div>{sub}'


def insight_card(icon, title, value, sub="", color="#7C3AED"):
    return f"""<div class="insight-card" style="--ic-color:{color}">
  <div class="ic-icon">{icon}</div>
  <div class="ic-title">{title}</div>
  <div class="ic-value">{value}</div>
  {f'<div class="ic-sub">{sub}</div>' if sub else ''}
</div>"""


def chart(fig, height=380):
    fig.update_layout(**CHART_LAYOUT, height=height, dragmode=False)
    st.plotly_chart(fig, use_container_width=True,
                    config={"displayModeBar": False, "scrollZoom": False,
                            "doubleClick": False, "staticPlot": False})


# ═══════════════════════════════════════════════════════════
# LOAD DATA
# ═══════════════════════════════════════════════════════════
orders_df, items_df, processing_df = load_data()
insights = compute_insights(orders_df, items_df)

# ═══════════════════════════════════════════════════════════
# SIDEBAR
# ═══════════════════════════════════════════════════════════
with st.sidebar:
    st.markdown("""
    <div class="sb-logo">
      <div style="font-size:28px; margin-bottom:8px;">⚡</div>
      <div class="sb-logo-title">MS ENTERPRISES</div>
      <div class="sb-logo-sub">Purchase Order Intelligence</div>
      <div class="sb-logo-badge">NAGPUR · MAHARASHTRA</div>
    </div>
    """, unsafe_allow_html=True)

    st.markdown('<div class="sb-nav-label">Navigation</div>', unsafe_allow_html=True)

    NAV_ITEMS = [
        "📊  Overview",
        "🗺️  Geography",
        "🧱  Materials",
        "👥  Clients",
        "📅  Time Trends",
        "💰  Financial",
        "📝  Data Entry",
        "⚙️  Processing",
    ]
    page = st.radio("nav", NAV_ITEMS, label_visibility="collapsed",
                    index=NAV_ITEMS.index(st.session_state.page)
                    if st.session_state.page in NAV_ITEMS else 0)
    st.session_state.page = page

    st.markdown('<div class="sb-nav-label" style="margin-top:16px;">Filters</div>', unsafe_allow_html=True)

    fy_opts = sorted(orders_df["financial_year"].dropna().unique().tolist())
    sel_fy  = st.multiselect("Financial Year", fy_opts, default=fy_opts,
                              label_visibility="collapsed",
                              placeholder="📅 All Financial Years")

    state_opts = sorted(orders_df["consignee_state"].dropna().unique().tolist())
    sel_states = st.multiselect("State", state_opts,
                                 label_visibility="collapsed",
                                 placeholder="🗺️ All States")

    mat_opts = sorted(items_df["material_category"].dropna().unique().tolist())
    sel_mats = st.multiselect("Category", mat_opts,
                               label_visibility="collapsed",
                               placeholder="🧱 All Categories")

    st.markdown("---")
    if st.button("🔄 Refresh Data", use_container_width=True):
        st.cache_data.clear(); st.rerun()

    # ── Full detail Excel export ──────────────────────────────
    try:
        _conn_ex = sqlite3.connect(DB_PATH)
        full_export = pd.read_sql_query("""
            SELECT
                o.po_number                     AS "Invoice / PO No.",
                o.document_date                 AS "Date",
                o.financial_year                AS "Financial Year",
                o.customer_name                 AS "Client / Factory",
                o.consignee_name                AS "Consignee Name",
                o.consignee_state               AS "State",
                o.consignee_city                AS "City",
                m.name_standardized             AS "Material",
                m.category                      AS "Category",
                ROUND(oi.quantity, 3)            AS "Quantity (MT)",
                ROUND(oi.rate, 2)               AS "Rate per MT (₹)",
                ROUND(oi.amount, 2)             AS "Line Amount (₹)",
                o.total_amount                  AS "PO Total Amount (₹)",
                o.vendor_name                   AS "Vendor",
                o.gstin                         AS "GST No."
            FROM orders o
            LEFT JOIN order_items oi ON oi.order_id = o.id
            LEFT JOIN materials   m  ON m.id = oi.material_id
            ORDER BY o.document_date DESC, o.id DESC
        """, _conn_ex)
        _conn_ex.close()

        # KG → MT normalisation for export
        buf_ex = io.BytesIO()
        with pd.ExcelWriter(buf_ex, engine="xlsxwriter") as writer:
            full_export.to_excel(writer, index=False, sheet_name="All Orders")
            wb  = writer.book
            ws  = writer.sheets["All Orders"]

            # Formatting
            hdr_fmt = wb.add_format({
                "bold": True, "bg_color": "#6366F1", "font_color": "#FFFFFF",
                "border": 1, "border_color": "#E2E8F0", "text_wrap": True,
                "valign": "vcenter", "align": "center"
            })
            num_fmt  = wb.add_format({"num_format": "#,##0.00", "border": 1, "border_color": "#E2E8F0"})
            text_fmt = wb.add_format({"border": 1, "border_color": "#E2E8F0"})
            alt_fmt  = wb.add_format({"bg_color": "#F1F5F9", "border": 1, "border_color": "#E2E8F0"})

            # Column widths
            col_widths = [18, 12, 12, 28, 28, 16, 16, 28, 18, 14, 16, 16, 18, 20, 18]
            for i, w in enumerate(col_widths):
                ws.set_column(i, i, w)
            ws.set_row(0, 28, hdr_fmt)

            # Re-write header with format
            for ci, col in enumerate(full_export.columns):
                ws.write(0, ci, col, hdr_fmt)

            # Data rows with alternating colour
            num_cols = {"Quantity (MT)", "Rate per MT (₹)", "Line Amount (₹)", "PO Total Amount (₹)"}
            for ri, row in enumerate(full_export.itertuples(index=False), start=1):
                bg = wb.add_format({
                    "bg_color": "#FFFFFF" if ri % 2 == 1 else "#F8FAFC",
                    "border": 1, "border_color": "#E2E8F0"
                })
                num_bg = wb.add_format({
                    "bg_color": "#FFFFFF" if ri % 2 == 1 else "#F8FAFC",
                    "num_format": "#,##0.00", "border": 1, "border_color": "#E2E8F0"
                })
                for ci, col in enumerate(full_export.columns):
                    val = getattr(row, col.replace(" ", "_").replace("/", "_").replace("(", "").replace(")", "").replace(".", "").replace("₹", ""), None)
                    val = row[ci]
                    ws.write(ri, ci, val, num_bg if col in num_cols else bg)

            # Freeze top row
            ws.freeze_panes(1, 0)
            ws.autofilter(0, 0, len(full_export), len(full_export.columns) - 1)

            # Summary sheet
            summary_data = {
                "Metric": ["Total Revenue (₹)", "Total Quantity (MT)", "Total Orders",
                           "Unique Clients", "Unique Materials", "States Covered"],
                "Value": [
                    full_export["PO Total Amount (₹)"].sum() / len(full_export["Invoice / PO No."].unique()) if not full_export.empty else 0,
                    full_export["Quantity (MT)"].sum(),
                    full_export["Invoice / PO No."].nunique(),
                    full_export["Client / Factory"].nunique(),
                    full_export["Material"].nunique(),
                    full_export["State"].nunique(),
                ]
            }
            # fix revenue
            summary_data["Value"][0] = orders_df["total_amount"].sum()
            pd.DataFrame(summary_data).to_excel(writer, index=False, sheet_name="Summary")

        st.download_button("📥 Download Full Excel", buf_ex.getvalue(),
                           "ms_enterprises_full_report.xlsx",
                           "application/vnd.ms-excel",
                           use_container_width=True)
    except Exception as e:
        buf = io.BytesIO()
        export_df = orders_df.copy()
        export_df["document_date"] = export_df["document_date"].astype(str)
        export_df.to_excel(buf, index=False, engine="xlsxwriter")
        st.download_button("📥 Download Report", buf.getvalue(),
                           "ms_enterprises_report.xlsx",
                           "application/vnd.ms-excel",
                           use_container_width=True)

    # PPT download
    if _PPTX_OK:
        ppt_bytes = generate_ppt(orders_df, items_df)
        if ppt_bytes:
            st.download_button(
                "📊 Download PPT",
                ppt_bytes,
                "ms_enterprises_analytics.pptx",
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
            )

    st.markdown(f"""
    <div class="sb-footer">
      <strong>Last updated</strong><br>
      {datetime.now().strftime("%d %b %Y · %H:%M")}
    </div>
    """, unsafe_allow_html=True)


# ═══════════════════════════════════════════════════════════
# APPLY FILTERS
# ═══════════════════════════════════════════════════════════
def apply_filters(df):
    d = df.copy()
    if sel_fy:     d = d[d["financial_year"].isin(sel_fy)]
    if sel_states: d = d[d["consignee_state"].isin(sel_states)]
    return d

fo = apply_filters(orders_df)  # filtered orders
fi = apply_filters(items_df)   # filtered items
if sel_mats:
    fi = fi[fi["material_category"].isin(sel_mats)]


# ═══════════════════════════════════════════════════════════
# COMPUTE KPI HELPERS
# ═══════════════════════════════════════════════════════════
def fy_growth(df, col="total_amount"):
    fy = df.groupby("financial_year")[col].sum().sort_index()
    if len(fy) < 2: return 0.0
    a, b = fy.iloc[-1], fy.iloc[-2]
    return (a - b) / b * 100 if b else 0.0

def monthly_vals(df, col="total_amount"):
    if df.empty or "document_date" not in df: return [0]*12
    d = df.dropna(subset=["document_date"]).copy()
    d["m"] = d["document_date"].dt.month
    m = d.groupby("m")[col].sum().reindex(range(1,13), fill_value=0)
    return m.tolist()


# ═══════════════════════════════════════════════════════════
# PAGE: OVERVIEW
# ═══════════════════════════════════════════════════════════
if page == "📊  Overview":
    st.markdown('<div class="page-title">Material <span>Dispatch Register</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="page-sub">What was supplied · to which factory · how much quantity · at what rate</div>', unsafe_allow_html=True)

    # ── KPI strip ────────────────────────────────────────────
    total_rev  = fo["total_amount"].sum()
    total_qty  = fi["quantity"].sum()
    n_cust     = fo["customer_name"].nunique()
    n_states   = fo["consignee_state"].nunique()
    rev_growth = fy_growth(fo, "total_amount")
    ord_growth = fy_growth(fo.assign(c=1), "c") if "financial_year" in fo else 0
    spark_rev  = monthly_vals(fo)
    spark_ord  = monthly_vals(fo.assign(total_amount=1))
    spark_qty  = monthly_vals(fi, "quantity") if not fi.empty else [0]*12

    render_kpi_hero([
        ("💰", "Total Revenue",    int(total_rev), fmt(total_rev),
         rev_growth, spark_rev, "#7C3AED","#5B21B6","rgba(124,58,237,0.25)", "currency"),
        ("📦", "Total Orders",     len(fo),        str(len(fo)),
         ord_growth, spark_ord, "#0891B2","#0E7490","rgba(8,145,178,0.25)", "count"),
        ("⚖️", "Total Qty (MT)",   int(total_qty), f"{total_qty:,.1f} MT",
         0, spark_qty, "#EA580C","#C2410C","rgba(234,88,12,0.25)", "qty"),
        ("🏭", "Factories Served", n_cust,         str(n_cust),
         0, spark_rev, "#059669","#047857","rgba(5,150,105,0.25)", "count"),
        ("🗺️", "States Covered",  n_states,       str(n_states),
         0, spark_rev, "#DB2777","#BE185D","rgba(219,39,119,0.25)", "count"),
    ])

    # ── MAIN: Material → Factory SANKEY flow diagram ───────────
    st.markdown(section_header("Material Flow — Where Did Each Material Go?",
        "Top 10 materials → Top 10 factories · link width = revenue"),
        unsafe_allow_html=True)

    if not fi.empty:
        top_m10  = fi.groupby("material_name")["amount"].sum().nlargest(10).index.tolist()
        top_f10  = fi.groupby("customer_name")["amount"].sum().nlargest(10).index.tolist()
        top_f10s = [f[:30] for f in top_f10]           # display names
        fac_map  = {f: s for f, s in zip(top_f10, top_f10s)}

        sk_data = (fi[fi["material_name"].isin(top_m10) & fi["customer_name"].isin(top_f10)]
                   .groupby(["material_name","customer_name"])
                   .agg(revenue=("amount","sum"), qty=("quantity","sum"),
                        avg_rate=("rate","mean"))
                   .reset_index())

        n_mat   = len(top_m10)
        mat_idx = {m: i        for i, m in enumerate(top_m10)}
        fac_idx = {f: n_mat+i  for i, f in enumerate(top_f10)}

        def _hex_rgba(h, a=0.35):
            h = h.lstrip("#")
            r,g,b = int(h[0:2],16),int(h[2:4],16),int(h[4:6],16)
            return f"rgba({r},{g},{b},{a})"

        sources, targets, values, qtys, rates_sk = [], [], [], [], []
        link_colors = []
        for _, row in sk_data.iterrows():
            if row["material_name"] in mat_idx and row["customer_name"] in fac_idx:
                si = mat_idx[row["material_name"]]
                sources.append(si)
                targets.append(fac_idx[row["customer_name"]])
                values.append(float(row["revenue"]))
                qtys.append(float(row["qty"]))
                rates_sk.append(float(row["avg_rate"]))
                link_colors.append(_hex_rgba(ACCENT[si % len(ACCENT)], 0.35))

        if dark_mode:
            node_fac_color = "rgba(34,211,238,0.25)"
            node_line      = "rgba(255,255,255,0.15)"
        else:
            node_fac_color = "rgba(99,102,241,0.12)"
            node_line      = "rgba(15,23,42,0.12)"

        node_colors = ([_hex_rgba(c, 0.85) for c in ACCENT[:n_mat]] +
                       [node_fac_color] * len(top_f10))
        all_labels  = top_m10 + top_f10s

        fig_sk = go.Figure(go.Sankey(
            arrangement="snap",
            node=dict(
                pad=18, thickness=22,
                line=dict(color=node_line, width=0.8),
                label=all_labels,
                color=node_colors,
                hovertemplate="<b>%{label}</b><br>Total: ₹%{value:,.0f}<extra></extra>",
            ),
            link=dict(
                source=sources, target=targets,
                value=values,
                color=link_colors,
                customdata=list(zip(qtys, rates_sk)),
                hovertemplate=(
                    "<b>%{source.label}</b> → <b>%{target.label}</b><br>"
                    "Revenue: ₹%{value:,.0f}<br>"
                    "Qty: %{customdata[0]:,.1f} MT<br>"
                    "Avg Rate: ₹%{customdata[1]:,.0f}/MT<extra></extra>"
                ),
            ),
        ))
        fig_sk.update_layout(**CHART_LAYOUT,
                             title="Material Flow: Top 10 Materials → Top 10 Factories",
                             height=500)
        fig_sk.update_layout(font=dict(size=11, family="Inter, sans-serif", color=_txt))
        chart(fig_sk, 500)

    # ── MATERIAL QUICK-LOOK ─────────────────────────────────
    st.markdown(section_header("Material Quick-Look",
        "Pick any material — see exactly where it went, how much qty, at what rate"),
        unsafe_allow_html=True)

    if not fi.empty:
        all_mats_sorted = (fi.groupby("material_name")["amount"]
                             .sum().sort_values(ascending=False).index.tolist())
        ql_mat = st.selectbox("Select material", all_mats_sorted,
                               label_visibility="collapsed",
                               placeholder="Choose material…")
        ql = fi[fi["material_name"] == ql_mat].copy()

        if not ql.empty:
            ql_rev  = ql["amount"].sum()
            ql_qty  = ql["quantity"].sum()
            ql_rate = ql["rate"].mean()
            ql_facs = ql["customer_name"].nunique()
            ql_pos  = ql["order_id"].nunique()

            mk1, mk2, mk3, mk4, mk5 = st.columns(5)
            mk1.metric("Total Revenue",  fmt(ql_rev))
            mk2.metric("Total Qty (MT)", f"{ql_qty:,.1f}")
            mk3.metric("Avg Rate / MT",  f"₹{ql_rate:,.0f}")
            mk4.metric("Factories",      str(ql_facs))
            mk5.metric("Total POs",      str(ql_pos))

            ql_c1, ql_c2 = st.columns([3, 2])
            with ql_c1:
                ql_fy = (ql.groupby(["customer_name","financial_year"])
                           .agg(qty=("quantity","sum"), revenue=("amount","sum"),
                                avg_rate=("rate","mean"), n_pos=("order_id","nunique"))
                           .reset_index())
                ql_fy["cust_short"] = ql_fy["customer_name"].str[:28]
                fac_order = (ql_fy.groupby("cust_short")["qty"]
                                  .sum().sort_values(ascending=True).index.tolist())
                years_u = sorted(ql_fy["financial_year"].dropna().unique())
                fig_ql = go.Figure()
                for i, yr in enumerate(years_u):
                    sub = (ql_fy[ql_fy["financial_year"]==yr]
                           .set_index("cust_short").reindex(fac_order).reset_index())
                    fig_ql.add_trace(go.Bar(
                        name=str(yr), x=sub["qty"].fillna(0), y=sub["cust_short"],
                        orientation="h",
                        marker=dict(color=ACCENT[i % len(ACCENT)], line=dict(width=0)),
                        customdata=list(zip(sub["revenue"].fillna(0),
                                           sub["avg_rate"].fillna(0),
                                           sub["n_pos"].fillna(0))),
                        hovertemplate=(
                            "<b>%{y}</b> · %{name}<br>"
                            "Qty: <b>%{x:,.1f} MT</b><br>"
                            "Revenue: ₹%{customdata[0]:,.0f}<br>"
                            "Avg Rate: ₹%{customdata[1]:,.0f}/MT<br>"
                            "POs: %{customdata[2]:.0f}<extra></extra>"
                        ),
                    ))
                ht = max(320, len(fac_order) * 38 + 80)
                fig_ql.update_layout(**CHART_LAYOUT,
                                     title=f"{ql_mat[:38]} — Qty by Factory & Year",
                                     height=ht, barmode="group")
                fig_ql.update_layout(
                    yaxis=dict(categoryorder="array", categoryarray=fac_order,
                               **CHART_LAYOUT["yaxis"]))
                chart(fig_ql, ht)

            with ql_c2:
                rate_yr = (ql.groupby("financial_year")
                             .agg(avg_rate=("rate","mean"), total_qty=("quantity","sum"),
                                  total_rev=("amount","sum"))
                             .reset_index().sort_values("financial_year"))
                fig_rt = go.Figure(go.Scatter(
                    x=rate_yr["financial_year"], y=rate_yr["avg_rate"],
                    mode="lines+markers+text",
                    line=dict(color=ACCENT[0], width=2.5, shape="spline"),
                    marker=dict(color=ACCENT[0], size=9,
                                line=dict(color="white", width=2)),
                    text=[f"₹{v:,.0f}" for v in rate_yr["avg_rate"]],
                    textposition="top center",
                    textfont=dict(size=10, color="#475569"),
                    customdata=list(zip(rate_yr["total_qty"], rate_yr["total_rev"])),
                    hovertemplate=(
                        "<b>%{x}</b><br>"
                        "Avg Rate: <b>₹%{y:,.0f}/MT</b><br>"
                        "Qty: %{customdata[0]:,.1f} MT<br>"
                        "Revenue: ₹%{customdata[1]:,.0f}<extra></extra>"
                    ),
                    fill="tozeroy", fillcolor="rgba(99,102,241,0.06)",
                ))
                fig_rt.update_layout(**CHART_LAYOUT,
                                     title="Avg Rate / MT by Year", height=200)
                chart(fig_rt, 200)

                fac_summ = (ql.groupby("customer_name")
                              .agg(qty=("quantity","sum"), revenue=("amount","sum"),
                                   avg_rate=("rate","mean"), n_pos=("order_id","nunique"),
                                   state=("consignee_state","first"))
                              .reset_index().sort_values("qty", ascending=False).head(8))
                fac_summ.columns = ["Factory","Qty (MT)","Revenue","Rate/MT","POs","State"]
                fac_summ["Qty (MT)"] = fac_summ["Qty (MT)"].apply(lambda x: f"{x:,.1f}")
                fac_summ["Revenue"]  = fac_summ["Revenue"].apply(fmt)
                fac_summ["Rate/MT"]  = fac_summ["Rate/MT"].apply(lambda x: f"₹{x:,.0f}")
                st.dataframe(fac_summ[["Factory","Qty (MT)","Revenue","Rate/MT","POs","State"]],
                             use_container_width=True, hide_index=True, height=260)

    # ── Year Revenue + Top Materials ─────────────────────────
    st.markdown(section_header("Year-wise Revenue"), unsafe_allow_html=True)
    bc1, bc2 = st.columns(2)

    with bc1:
        fy_data = fo.groupby("financial_year").agg(
            revenue=("total_amount","sum"), orders=("id","count")
        ).reset_index().sort_values("financial_year")
        if not fy_data.empty:
            qty_by_fy = (fi.groupby("financial_year")["quantity"].sum()
                           .reindex(fy_data["financial_year"]).fillna(0).tolist())
            avg_ord   = (fy_data["revenue"] / fy_data["orders"]).tolist()
            fig_fy = go.Figure(go.Bar(
                x=fy_data["financial_year"], y=fy_data["revenue"],
                marker=dict(color=ACCENT[:len(fy_data)], line=dict(width=0)),
                text=[fmt(v) for v in fy_data["revenue"]],
                textposition="outside",
                textfont=dict(color="#64748B", size=11),
                customdata=list(zip(fy_data["orders"], qty_by_fy, avg_ord)),
                hovertemplate=(
                    "<b>%{x}</b><br>"
                    "Revenue: <b>%{text}</b><br>"
                    "Orders: %{customdata[0]}<br>"
                    "Total Qty: %{customdata[1]:,.1f} MT<br>"
                    "Avg Order Value: ₹%{customdata[2]:,.0f}<extra></extra>"
                ),
            ))
            fig_fy.update_layout(**CHART_LAYOUT, title="Revenue by Financial Year", height=280)
            chart(fig_fy, 280)

    with bc2:
        if not fi.empty:
            top_m_agg = (fi.groupby("material_name")
                           .agg(revenue=("amount","sum"), qty=("quantity","sum"),
                                avg_rate=("rate","mean"),
                                n_fac=("customer_name","nunique"),
                                n_pos=("order_id","nunique"))
                           .nlargest(8,"revenue").reset_index().sort_values("revenue"))
            fig_tm = go.Figure(go.Bar(
                x=top_m_agg["revenue"], y=top_m_agg["material_name"].str[:30],
                orientation="h",
                marker=dict(color=top_m_agg["revenue"],
                            colorscale=[[0,"#C4B5FD"],[1,"#6366F1"]],
                            showscale=False, line=dict(width=0)),
                text=[fmt(v) for v in top_m_agg["revenue"]],
                textposition="outside",
                textfont=dict(color="#64748B", size=10),
                customdata=list(zip(top_m_agg["qty"], top_m_agg["avg_rate"],
                                    top_m_agg["n_fac"], top_m_agg["n_pos"])),
                hovertemplate=(
                    "<b>%{y}</b><br>"
                    "Revenue: <b>%{text}</b><br>"
                    "Total Qty: %{customdata[0]:,.1f} MT<br>"
                    "Avg Rate: ₹%{customdata[1]:,.0f}/MT<br>"
                    "Factories supplied: %{customdata[2]}<br>"
                    "Total POs: %{customdata[3]}<extra></extra>"
                ),
            ))
            fig_tm.update_layout(**CHART_LAYOUT,
                                 title="Top 8 Materials — hover for full detail", height=280)
            chart(fig_tm, 280)


# ═══════════════════════════════════════════════════════════
# PAGE: GEOGRAPHY
# ═══════════════════════════════════════════════════════════
elif page == "🗺️  Geography":
    st.markdown('<div class="page-title">Geographic <span>Analysis</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="page-sub">State-wise supply distribution and regional revenue analysis</div>', unsafe_allow_html=True)

    state_rev = fo.groupby("consignee_state").agg(
        revenue=("total_amount","sum"), orders=("id","count")).reset_index().sort_values("revenue", ascending=False)

    # Top state metrics
    top4 = state_rev.head(4)
    total_r = fo["total_amount"].sum() or 1
    render_kpi_hero([
        (str(r["consignee_state"])[:3].upper(),
         str(r["consignee_state"]),
         int(r["revenue"]), fmt(r["revenue"]),
         0, [],
         ACCENT[i], ACCENT[i], ACCENT[i]+"44")
        for i, (_, r) in enumerate(top4.iterrows())
    ] + [("","","","",0,[],ACCENT[4],ACCENT[4],ACCENT[4]+"11")] * (5 - len(top4)))

    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2 = st.columns(2)

    with c1:
        st.markdown(section_header("Revenue by State"), unsafe_allow_html=True)
        if not state_rev.empty:
            fig = go.Figure(go.Bar(
                x=state_rev["consignee_state"], y=state_rev["revenue"],
                marker=dict(
                    color=state_rev["revenue"],
                    colorscale=[[0,"#3B1F7A"],[0.5,"#7C3AED"],[1,"#C4B5FD"]],
                    showscale=False, line=dict(width=0),
                ),
                text=[fmt(v) for v in state_rev["revenue"]],
                textposition="outside",
                hovertemplate="<b>%{x}</b><br>Revenue: %{text}<extra></extra>",
            ))
            fig.update_layout(**CHART_LAYOUT, title="State Revenue (Sorted)", height=350)
            chart(fig, 350)

    with c2:
        st.markdown(section_header("Revenue Share — State Donut"), unsafe_allow_html=True)
        if not state_rev.empty:
            fig2 = go.Figure(go.Pie(
                labels=state_rev["consignee_state"],
                values=state_rev["revenue"],
                hole=0.55,
                marker=dict(colors=ACCENT[:len(state_rev)],
                            line=dict(color="#0B0F1A", width=2)),
                textinfo="percent",
                textfont=dict(size=11, color="#E5E7EB"),
                hovertemplate="<b>%{label}</b><br>Revenue: ₹%{value:,.0f}<br>Share: %{percent}<extra></extra>",
            ))
            fig2.add_annotation(text=f"{len(state_rev)}<br>States", x=0.5, y=0.5,
                                font=dict(size=16, color="#E5E7EB", family="Inter"),
                                showarrow=False)
            fig2.update_layout(**CHART_LAYOUT, title="Revenue Distribution", height=350)
            chart(fig2, 350)

    # India Map
    st.markdown(section_header("India Supply Heatmap"), unsafe_allow_html=True)
    try:
        import folium
        from streamlit_folium import st_folium

        STATE_COORDS = {
            "Maharashtra":[19.7515,75.7139],"Jharkhand":[23.6102,85.2799],
            "Chhattisgarh":[21.2787,81.8661],"Madhya Pradesh":[22.9734,78.6569],
            "Odisha":[20.9517,85.0985],"Andhra Pradesh":[15.9129,79.7400],
            "Rajasthan":[27.0238,74.2179],"West Bengal":[22.9868,87.8550],
            "Uttar Pradesh":[26.8467,80.9462],"Telangana":[18.1124,79.0193],
            "Karnataka":[15.3173,75.7139],"Tamil Nadu":[11.1271,78.6569],
            "Gujarat":[22.2587,71.1924],"Delhi":[28.7041,77.1025],
        }
        m = folium.Map(location=[22.5,80], zoom_start=5, tiles="CartoDB dark_matter")
        max_rev = state_rev["revenue"].max() or 1
        for _, r in state_rev.iterrows():
            if r["consignee_state"] in STATE_COORDS:
                coord = STATE_COORDS[r["consignee_state"]]
                radius = 15 + (r["revenue"] / max_rev) * 45
                folium.CircleMarker(
                    location=coord, radius=radius,
                    color="#7C3AED", fill=True, fill_color="#7C3AED", fill_opacity=0.6,
                    popup=folium.Popup(
                        f"<div style='background:#1a2035;padding:10px;border-radius:8px;color:#E5E7EB'>"
                        f"<b style='color:#A78BFA'>{r['consignee_state']}</b><br>"
                        f"Revenue: ₹{r['revenue']:,.0f}<br>Orders: {r['orders']}</div>",
                        max_width=200),
                    tooltip=f"{r['consignee_state']}: {r['orders']} orders"
                ).add_to(m)
        st_folium(m, width=None, height=440, returned_objects=[])
    except ImportError:
        st.info("Run `pip install folium streamlit-folium` for interactive map")

    # State × Year heatmap
    st.markdown(section_header("State Revenue by Year"), unsafe_allow_html=True)
    sy = fo.groupby(["consignee_state","financial_year"])["total_amount"].sum().reset_index()
    if not sy.empty:
        pivot = sy.pivot_table(index="consignee_state", columns="financial_year",
                               values="total_amount", fill_value=0)
        fig3 = go.Figure(go.Heatmap(
            z=pivot.values, x=pivot.columns.tolist(), y=pivot.index.tolist(),
            colorscale=[[0,"#0B0F1A"],[0.3,"#3B1F7A"],[0.7,"#7C3AED"],[1,"#C4B5FD"]],
            text=[[fmt(v) for v in row] for row in pivot.values],
            texttemplate="%{text}", textfont=dict(size=11),
            hovertemplate="<b>%{y}</b> — %{x}<br>Revenue: %{text}<extra></extra>",
        ))
        fig3.update_layout(**CHART_LAYOUT, title="Revenue Heatmap: State × Year", height=300)
        chart(fig3, 300)


# ═══════════════════════════════════════════════════════════
# PAGE: MATERIALS
# ═══════════════════════════════════════════════════════════
elif page == "🧱  Materials":
    st.markdown('<div class="page-title">Material <span>Flow</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="page-sub">What was supplied, where it went, in what quantity, at what price</div>', unsafe_allow_html=True)

    mat_agg = fi.groupby("material_name").agg(
        revenue=("amount","sum"), qty=("quantity","sum"), orders=("order_id","nunique")
    ).reset_index().sort_values("revenue", ascending=False)
    cat_agg = fi.groupby("material_category")["amount"].sum().reset_index().sort_values("amount", ascending=False)
    cat_agg = cat_agg[cat_agg["amount"] > 0]
    total_item_rev = fi["amount"].sum() or 1
    top_mat = mat_agg.iloc[0]["material_name"] if not mat_agg.empty else "N/A"
    top_cat = cat_agg.iloc[0]["material_category"] if not cat_agg.empty else "N/A"
    n_mats  = mat_agg["material_name"].nunique()

    # KPI row
    mc1, mc2, mc3, mc4 = st.columns(4)
    mc1.metric("Total Material Revenue", fmt(total_item_rev))
    mc2.metric("Unique Materials", str(n_mats))
    mc3.metric("Top Material", top_mat[:22])
    mc4.metric("Top Category", top_cat)

    # ── MAIN CHART: Material → Factory treemap ──────────────
    st.markdown(section_header("Material Supply Overview",
                "Every material · every factory · sized by revenue"),
                unsafe_allow_html=True)
    if not fi.empty:
        tree_df = fi.groupby(["material_category","material_name","customer_name"]).agg(
            revenue=("amount","sum"), qty=("quantity","sum")
        ).reset_index()
        tree_df = tree_df[tree_df["revenue"] > 0]
        tree_df["cust_short"] = tree_df["customer_name"].str[:30]
        tree_df["hover"] = (tree_df["material_name"] + "<br>→ " +
                            tree_df["cust_short"] + "<br>Qty: " +
                            tree_df["qty"].apply(lambda x: f"{x:,.1f}") +
                            " MT  |  " + tree_df["revenue"].apply(fmt))
        fig_tree = go.Figure(go.Treemap(
            ids   = (tree_df["material_category"] + "/" +
                     tree_df["material_name"] + "/" + tree_df["cust_short"]),
            labels= tree_df["cust_short"],
            parents=(tree_df["material_category"] + "/" + tree_df["material_name"]),
            values= tree_df["revenue"],
            customdata=tree_df[["material_name","qty","revenue","cust_short"]],
            hovertemplate=(
                "<b>%{customdata[0]}</b><br>"
                "Factory: %{customdata[3]}<br>"
                "Qty: %{customdata[1]:,.1f} MT<br>"
                "Revenue: ₹%{customdata[2]:,.0f}<extra></extra>"
            ),
            branchvalues="total",
            marker=dict(
                colorscale=[[0,"#EDE9FE"],[0.4,"#7C3AED"],[1,"#4C1D95"]],
                cmid=total_item_rev / (len(tree_df)+1),
                showscale=False,
                line=dict(color="white", width=2),
            ),
            textfont=dict(size=11, family="Inter"),
        ))
        # Add parent (material) level
        mat_tree = fi.groupby(["material_category","material_name"]).agg(
            revenue=("amount","sum")).reset_index()
        cat_tree = fi.groupby("material_category").agg(revenue=("amount","sum")).reset_index()
        fig_tree2 = go.Figure(go.Treemap(
            ids    = ([""] +
                      cat_tree["material_category"].tolist() +
                      (mat_tree["material_category"]+"/"+mat_tree["material_name"]).tolist() +
                      (tree_df["material_category"]+"/"+tree_df["material_name"]+"/"+tree_df["cust_short"]).tolist()),
            labels = (["All Materials"] +
                      cat_tree["material_category"].tolist() +
                      mat_tree["material_name"].tolist() +
                      tree_df["cust_short"].tolist()),
            parents= ([""] +
                      [""]*len(cat_tree) +
                      cat_tree["material_category"].tolist() +  # This won't work directly, need to map
                      (tree_df["material_category"]+"/"+tree_df["material_name"]).tolist()),
            values = ([0] +
                      cat_tree["revenue"].tolist() +
                      mat_tree["revenue"].tolist() +
                      tree_df["revenue"].tolist()),
            branchvalues="remainder",
            marker=dict(
                colors=(["#F1F5F9"] +
                        ACCENT[:len(cat_tree)] +
                        [ACCENT[i % len(ACCENT)] for i in range(len(mat_tree))] +
                        [ACCENT[i % len(ACCENT)] for i in range(len(tree_df))]),
                line=dict(color="white", width=2),
            ),
            textfont=dict(size=11, family="Inter"),
            hovertemplate="<b>%{label}</b><br>Revenue: ₹%{value:,.0f}<extra></extra>",
        ))
        # Fix parent mapping for mat_tree
        mat_parent_map = mat_tree.set_index("material_name")["material_category"].to_dict()
        fig_tree2.data[0].parents = (
            [""] +
            [""]*len(cat_tree) +
            [mat_tree.iloc[i]["material_category"] for i in range(len(mat_tree))] +
            (tree_df["material_category"]+"/"+tree_df["material_name"]).tolist()
        )
        fig_tree2.update_layout(**CHART_LAYOUT, title="Material → Factory Revenue Map", height=420)
        fig_tree2.update_layout(margin=dict(l=0, r=0, t=44, b=0))
        chart(fig_tree2, 420)

    # ── MATERIAL → FACTORY FLOW: Stacked bar ───────────────
    st.markdown(section_header("Where Did Each Material Go?",
                "Top 10 materials · stacked by destination factory"),
                unsafe_allow_html=True)
    if not fi.empty:
        top10_mats = mat_agg.head(10)["material_name"].tolist()
        flow = fi[fi["material_name"].isin(top10_mats)].copy()
        flow["cust_short"] = flow["customer_name"].str[:28]
        flow_agg = flow.groupby(["material_name","cust_short"])["amount"].sum().reset_index()
        # Build stacked horizontal bar
        factories = flow_agg["cust_short"].unique().tolist()
        fig_flow = go.Figure()
        colors_f = ACCENT * 4
        for i, fac in enumerate(factories):
            sub = flow_agg[flow_agg["cust_short"]==fac]
            mat_rev = []
            for m in top10_mats:
                row = sub[sub["material_name"]==m]
                mat_rev.append(float(row["amount"].values[0]) if not row.empty else 0)
            fig_flow.add_trace(go.Bar(
                name=fac, x=mat_rev,
                y=[m[:30] for m in top10_mats],
                orientation="h",
                marker=dict(color=colors_f[i], line=dict(width=0)),
                hovertemplate="<b>%{y}</b><br>→ " + fac + "<br>₹%{x:,.0f}<extra></extra>",
            ))
        fig_flow.update_layout(**CHART_LAYOUT, title="Material Flow to Factories (Revenue)",
                               height=420, barmode="stack")
        fig_flow.update_layout(yaxis=dict(categoryorder="total ascending", **CHART_LAYOUT["yaxis"]))
        chart(fig_flow, 420)

    # ── DRILL-DOWN: Select a material, see full details ─────
    st.markdown(section_header("Material Drill-Down",
                "Select any material to see factory · quantity · year · rate"),
                unsafe_allow_html=True)
    mat_choices = mat_agg["material_name"].tolist()
    if mat_choices:
        sel_mat = st.selectbox("Select Material", mat_choices,
                               label_visibility="collapsed",
                               placeholder="Choose a material…")
        drill = fi[fi["material_name"] == sel_mat].copy()
        if not drill.empty:
            drill["year_label"] = drill["financial_year"].fillna(drill["year"].astype(str))
            drill_agg = drill.groupby(["customer_name","year_label"]).agg(
                qty=("quantity","sum"),
                total_rev=("amount","sum"),
                avg_rate=("rate","mean"),
                orders=("order_id","nunique"),
            ).reset_index().sort_values("total_rev", ascending=False)

            # Summary KPIs for this material
            dk1, dk2, dk3, dk4 = st.columns(4)
            dk1.metric("Total Revenue", fmt(drill["amount"].sum()))
            dk2.metric("Total Qty (MT)", f"{drill['quantity'].sum():,.1f}")
            dk3.metric("Avg Rate / MT", fmt(drill["rate"].mean()))
            dk4.metric("Factories Supplied", str(drill["customer_name"].nunique()))

            # Chart: qty by factory+year
            dc1, dc2 = st.columns([3,2])
            with dc1:
                drill_bar = drill.groupby(["customer_name","year_label"]).agg(
                    qty=("quantity","sum")).reset_index()
                drill_bar["cust_short"] = drill_bar["customer_name"].str[:28]
                years_u = sorted(drill_bar["year_label"].dropna().unique())
                fig_d = go.Figure()
                for i, yr in enumerate(years_u):
                    sub = drill_bar[drill_bar["year_label"]==yr]
                    fig_d.add_trace(go.Bar(
                        name=str(yr),
                        x=sub["qty"], y=sub["cust_short"],
                        orientation="h",
                        marker=dict(color=ACCENT[i % len(ACCENT)], line=dict(width=0)),
                        hovertemplate="<b>%{y}</b> (%{name})<br>Qty: %{x:,.1f} MT<extra></extra>",
                    ))
                fig_d.update_layout(**CHART_LAYOUT,
                                    title=f"Qty Supplied: {sel_mat[:40]} — by Factory & Year",
                                    height=max(300, len(drill_bar["cust_short"].unique())*45+80),
                                    barmode="group")
                fig_d.update_layout(yaxis=dict(categoryorder="total ascending",
                                               **CHART_LAYOUT["yaxis"]))
                chart(fig_d, max(300, len(drill_bar["cust_short"].unique())*45+80))

            with dc2:
                # Price trend by year
                price_yr = drill.groupby("year_label")["rate"].mean().reset_index().sort_values("year_label")
                if len(price_yr) > 1:
                    fig_p = go.Figure(go.Scatter(
                        x=price_yr["year_label"], y=price_yr["rate"],
                        mode="lines+markers+text",
                        line=dict(color="#7C3AED", width=2.5),
                        marker=dict(color="#7C3AED", size=8),
                        text=[fmt(v) for v in price_yr["rate"]],
                        textposition="top center",
                        textfont=dict(size=10, color="#475569"),
                        hovertemplate="<b>%{x}</b><br>Avg Rate: ₹%{y:,.0f}/MT<extra></extra>",
                    ))
                    fig_p.update_layout(**CHART_LAYOUT,
                                       title="Avg Rate per MT by Year", height=260)
                    chart(fig_p, 260)

            # Detail table
            display_cols = ["customer_name","year_label","qty","avg_rate","total_rev","orders"]
            rename_map = {"customer_name":"Factory","year_label":"FY",
                          "qty":"Qty (MT)","avg_rate":"Avg Rate/MT",
                          "total_rev":"Revenue","orders":"POs"}
            disp = drill_agg[display_cols].rename(columns=rename_map).copy()
            disp["Qty (MT)"]    = disp["Qty (MT)"].apply(lambda x: f"{x:,.1f}")
            disp["Avg Rate/MT"] = disp["Avg Rate/MT"].apply(lambda x: f"₹{x:,.0f}")
            disp["Revenue"]     = disp["Revenue"].apply(fmt)
            st.dataframe(disp, use_container_width=True, hide_index=True)

    # ── Category split + Year trend ─────────────────────────
    cc1, cc2 = st.columns([2,3])
    with cc1:
        st.markdown(section_header("Category Mix"), unsafe_allow_html=True)
        if not cat_agg.empty:
            fig_pie = go.Figure(go.Pie(
                labels=cat_agg["material_category"], values=cat_agg["amount"],
                hole=0.55,
                marker=dict(colors=ACCENT[:len(cat_agg)],
                            line=dict(color="white", width=2)),
                textinfo="percent+label",
                textfont=dict(size=11, color="#1e293b"),
                hovertemplate="<b>%{label}</b><br>₹%{value:,.0f}<br>%{percent}<extra></extra>",
            ))
            fig_pie.add_annotation(text="Category<br>Mix", x=0.5, y=0.5,
                                   font=dict(size=12, color="#64748B"), showarrow=False)
            fig_pie.update_layout(**CHART_LAYOUT, title="Revenue by Category", height=300)
            chart(fig_pie, 300)

    with cc2:
        st.markdown(section_header("Category Revenue by Year"), unsafe_allow_html=True)
        mat_year = fi.groupby(["financial_year","material_category"])["amount"].sum().reset_index()
        if not mat_year.empty:
            fig_yr = px.bar(mat_year, x="financial_year", y="amount",
                            color="material_category", barmode="stack",
                            color_discrete_sequence=ACCENT)
            fig_yr.update_layout(**CHART_LAYOUT, title="Category Trend by Year", height=300)
            chart(fig_yr, 300)


# ═══════════════════════════════════════════════════════════
# PAGE: CLIENTS
# ═══════════════════════════════════════════════════════════
elif page == "👥  Clients":
    st.markdown('<div class="page-title">Client <span>Performance</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="page-sub">Customer performance, retention and revenue contribution</div>', unsafe_allow_html=True)

    cust_agg = fo.groupby("customer_name").agg(
        revenue=("total_amount","sum"), orders=("id","count"),
        first=("document_date","min"), last=("document_date","max"),
    ).reset_index().sort_values("revenue", ascending=False)

    n_repeat  = len(cust_agg[cust_agg["orders"] > 1])
    n_onetime = len(cust_agg[cust_agg["orders"] == 1])
    retention = n_repeat / len(cust_agg) * 100 if len(cust_agg) else 0

    cc1, cc2, cc3, cc4 = st.columns(4)
    cc1.metric("Total Customers",   str(len(cust_agg)))
    cc2.metric("Repeat Customers",  str(n_repeat),  delta=f"{n_repeat} of {len(cust_agg)}")
    cc3.metric("One-time Clients",  str(n_onetime))
    cc4.metric("Retention Rate",    f"{retention:.0f}%")

    c1, c2 = st.columns([2, 3])

    with c1:
        st.markdown(section_header("Top Clients Leaderboard"), unsafe_allow_html=True)
        total_rev = cust_agg["revenue"].sum() or 1
        max_rev   = cust_agg["revenue"].max() or 1
        lb_html = ""
        for i, (_, r) in enumerate(cust_agg.head(10).iterrows()):
            rank_cls = ["gold","silver","bronze","","","","","","",""][i]
            pct = r["revenue"] / total_rev * 100
            bar_w = int(r["revenue"] / max_rev * 100)
            accent_c = ACCENT[i % len(ACCENT)]
            lb_html += f"""<div class="lb-item">
              <div class="lb-rank {rank_cls}">#{i+1}</div>
              <div style="flex:1;min-width:0">
                <div class="lb-name" style="font-size:12px;white-space:nowrap;overflow:hidden;text-overflow:ellipsis">{str(r['customer_name'])[:32]}</div>
                <div class="lb-sub">{r['orders']} orders</div>
                <div class="lb-bar-wrap"><div class="lb-bar" style="width:{bar_w}%;background:linear-gradient(90deg,{accent_c},{accent_c}88)"></div></div>
              </div>
              <div class="lb-rev">
                <div style="font-size:13px;font-weight:700;color:#F9FAFB">{fmt(r['revenue'])}</div>
                <div style="font-size:11px;color:#6B7280;text-align:right">{pct:.1f}%</div>
              </div>
            </div>"""
        st.markdown(lb_html, unsafe_allow_html=True)

    with c2:
        st.markdown(section_header("Revenue by Customer (Top 15)"), unsafe_allow_html=True)
        top15c = cust_agg.head(15)
        if not top15c.empty:
            fig = go.Figure(go.Bar(
                x=top15c["revenue"], y=top15c["customer_name"].str[:30],
                orientation="h",
                marker=dict(
                    color=top15c["revenue"],
                    colorscale=[[0,"#0C4A6E"],[0.5,"#06B6D4"],[1,"#67E8F9"]],
                    showscale=False, line=dict(width=0),
                ),
                text=[fmt(v) for v in top15c["revenue"]],
                textposition="outside",
                textfont=dict(color="#9CA3AF", size=10),
                hovertemplate="<b>%{y}</b><br>Revenue: %{text}<extra></extra>",
            ))
            fig.update_layout(**CHART_LAYOUT, title="Top 15 by Revenue", height=500)
            fig.update_layout(yaxis=dict(categoryorder="total ascending", **CHART_LAYOUT["yaxis"]))
            chart(fig, 500)

    # Client × Year heatmap
    st.markdown(section_header("Client Activity Heatmap"), unsafe_allow_html=True)
    cy = fo.groupby(["customer_name","financial_year"])["total_amount"].sum().reset_index()
    if not cy.empty:
        top_custs = cust_agg.head(15)["customer_name"].tolist()
        cy2 = cy[cy["customer_name"].isin(top_custs)]
        pivot = cy2.pivot_table(index="customer_name", columns="financial_year",
                                values="total_amount", fill_value=0)
        pivot.index = [str(x)[:28] for x in pivot.index]
        fig2 = go.Figure(go.Heatmap(
            z=pivot.values, x=pivot.columns.tolist(), y=pivot.index.tolist(),
            colorscale=[[0,"#0B0F1A"],[0.3,"#0C4A6E"],[0.7,"#06B6D4"],[1,"#BAE6FD"]],
            text=[[fmt(v) for v in row] for row in pivot.values],
            texttemplate="%{text}", textfont=dict(size=10),
            hovertemplate="<b>%{y}</b> — %{x}<br>%{text}<extra></extra>",
        ))
        fig2.update_layout(**CHART_LAYOUT, title="Top 15 Clients · Revenue by Year", height=420)
        chart(fig2, 420)


# ═══════════════════════════════════════════════════════════
# PAGE: TIME TRENDS
# ═══════════════════════════════════════════════════════════
elif page == "📅  Time Trends":
    st.markdown('<div class="page-title">Time <span>Trends</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="page-sub">Seasonality, year-over-year growth and cumulative performance</div>', unsafe_allow_html=True)

    tdf = fo.dropna(subset=["document_date"]).copy()
    tdf["period"] = tdf["document_date"].dt.to_period("M").astype(str)
    tdf["quarter"] = tdf["document_date"].dt.month.map(
        {1:"Q4",2:"Q4",3:"Q4",4:"Q1",5:"Q1",6:"Q1",7:"Q2",8:"Q2",9:"Q2",10:"Q3",11:"Q3",12:"Q3"})
    tdf["month_name"] = tdf["document_date"].dt.strftime("%b")
    tdf["month_num"]  = tdf["document_date"].dt.month

    c1, c2 = st.columns(2)

    with c1:
        st.markdown(section_header("Monthly Revenue Trend"), unsafe_allow_html=True)
        if not tdf.empty:
            m_data = tdf.groupby("period")["total_amount"].sum().reset_index().sort_values("period")
            fig = go.Figure()
            fig.add_trace(go.Scatter(
                x=m_data["period"], y=m_data["total_amount"],
                mode="lines+markers", name="Revenue",
                line=dict(color="#7C3AED", width=2.5, shape="spline"),
                marker=dict(color="#7C3AED", size=5, line=dict(color="#A78BFA",width=1.5)),
                fill="tozeroy",
                fillgradient=dict(colorscale=[[0,"rgba(124,58,237,0.3)"],[1,"rgba(124,58,237,0)"]]),
                hovertemplate="<b>%{x}</b><br>₹%{y:,.0f}<extra></extra>",
            ))
            fig.update_layout(**CHART_LAYOUT, title="Monthly Revenue (All Time)", height=300)
            fig.update_layout(xaxis=dict(tickangle=-45, **CHART_LAYOUT["xaxis"]))
            chart(fig, 300)

    with c2:
        st.markdown(section_header("Seasonality — Month Pattern"), unsafe_allow_html=True)
        if not tdf.empty:
            MN = {1:"Jan",2:"Feb",3:"Mar",4:"Apr",5:"May",6:"Jun",
                  7:"Jul",8:"Aug",9:"Sep",10:"Oct",11:"Nov",12:"Dec"}
            s_data = tdf.groupby("month_num")["total_amount"].sum().reset_index()
            s_data["month_name"] = s_data["month_num"].map(MN)
            s_data = s_data.sort_values("month_num")
            fig2 = go.Figure(go.Bar(
                x=s_data["month_name"], y=s_data["total_amount"],
                marker=dict(
                    color=s_data["total_amount"],
                    colorscale=[[0,"#0C4A6E"],[0.5,"#06B6D4"],[1,"#67E8F9"]],
                    showscale=False, line=dict(width=0),
                ),
                hovertemplate="<b>%{x}</b><br>₹%{y:,.0f}<extra></extra>",
            ))
            fig2.update_layout(**CHART_LAYOUT, title="Revenue by Month (Cumulative)", height=300)
            chart(fig2, 300)

    # YoY growth
    st.markdown(section_header("Year-over-Year Growth"), unsafe_allow_html=True)
    fy_rev = fo.groupby("financial_year")["total_amount"].sum().reset_index().sort_values("financial_year")
    fy_rev["growth"] = fy_rev["total_amount"].pct_change() * 100
    if not fy_rev.empty:
        fig3 = go.Figure()
        fig3.add_trace(go.Bar(
            x=fy_rev["financial_year"], y=fy_rev["total_amount"],
            name="Revenue", marker=dict(color=ACCENT[0], opacity=0.85, line=dict(width=0)),
            hovertemplate="<b>%{x}</b><br>Revenue: ₹%{y:,.0f}<extra></extra>",
        ))
        fig3.add_trace(go.Scatter(
            x=fy_rev["financial_year"], y=fy_rev["growth"],
            name="YoY Growth %", yaxis="y2",
            mode="lines+markers+text",
            line=dict(color="#10B981", width=2.5),
            marker=dict(color="#10B981", size=8, line=dict(color="#34D399", width=2)),
            text=[f"{g:.0f}%" if pd.notna(g) else "" for g in fy_rev["growth"]],
            textposition="top center", textfont=dict(color="#34D399", size=12),
            hovertemplate="<b>%{x}</b><br>Growth: %{y:.1f}%<extra></extra>",
        ))
        fig3.update_layout(**CHART_LAYOUT, title="Revenue & YoY Growth", height=320)
        fig3.update_layout(
            yaxis=dict(title="Revenue (₹)", gridcolor="rgba(255,255,255,0.05)",
                       tickfont=dict(color="#6B7280"), showline=False),
            yaxis2=dict(title="Growth %", overlaying="y", side="right",
                        showgrid=False, tickfont=dict(color="#10B981"),
                        zeroline=True, zerolinecolor="rgba(255,255,255,0.1)"),
        )
        chart(fig3, 320)

    # Cumulative
    st.markdown(section_header("Cumulative Revenue Growth"), unsafe_allow_html=True)
    cum_df = tdf.sort_values("document_date").copy()
    cum_df["cumulative"] = cum_df["total_amount"].cumsum()
    if not cum_df.empty:
        fig4 = go.Figure(go.Scatter(
            x=cum_df["document_date"], y=cum_df["cumulative"],
            mode="lines", name="Cumulative",
            line=dict(color="#F97316", width=2.5, shape="spline"),
            fill="tozeroy",
            fillgradient=dict(colorscale=[[0,"rgba(249,115,22,0.3)"],[1,"rgba(249,115,22,0)"]]),
            hovertemplate="<b>%{x|%d %b %Y}</b><br>Cumulative: ₹%{y:,.0f}<extra></extra>",
        ))
        fig4.update_layout(**CHART_LAYOUT, title="Cumulative Revenue Over Time", height=270)
        chart(fig4, 270)


# ═══════════════════════════════════════════════════════════
# PAGE: FINANCIAL
# ═══════════════════════════════════════════════════════════
elif page == "💰  Financial":
    st.markdown('<div class="page-title">Financial <span>Analytics</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="page-sub">Revenue, invoice values and order distribution analysis</div>', unsafe_allow_html=True)

    total_rev = fo["total_amount"].sum()
    avg_rev   = fo["total_amount"].mean() if len(fo) else 0
    max_rev   = fo["total_amount"].max() if len(fo) else 0
    gst_total = fo["gst_amount"].sum()

    render_kpi_hero([
        ("💰","Total Revenue",    int(total_rev), fmt(total_rev), 0, [], ACCENT[0], ACCENT[0], ACCENT[0]+"44"),
        ("📊","Avg Order Value",  int(avg_rev),   fmt(avg_rev),   0, [], ACCENT[1], ACCENT[1], ACCENT[1]+"44"),
        ("🏆","Largest Order",    int(max_rev),   fmt(max_rev),   0, [], ACCENT[2], ACCENT[2], ACCENT[2]+"44"),
        ("🧾","Total GST",        int(gst_total), fmt(gst_total), 0, [], ACCENT[3], ACCENT[3], ACCENT[3]+"44"),
    ])

    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2 = st.columns(2)

    with c1:
        st.markdown(section_header("Revenue Share by State"), unsafe_allow_html=True)
        sr = fo.groupby("consignee_state")["total_amount"].sum().reset_index().sort_values("total_amount", ascending=False)
        if not sr.empty:
            fig = go.Figure(go.Pie(
                labels=sr["consignee_state"], values=sr["total_amount"], hole=0.55,
                marker=dict(colors=ACCENT[:len(sr)], line=dict(color="#0B0F1A", width=2)),
                textinfo="percent", textfont=dict(size=11, color="#E5E7EB"),
                hovertemplate="<b>%{label}</b><br>₹%{value:,.0f}<br>%{percent}<extra></extra>",
            ))
            fig.add_annotation(text="State<br>Revenue", x=0.5, y=0.5,
                               font=dict(size=13, color="#9CA3AF"), showarrow=False)
            fig.update_layout(**CHART_LAYOUT, title="State Revenue Distribution", height=320)
            chart(fig, 320)

    with c2:
        st.markdown(section_header("Revenue Share by Customer"), unsafe_allow_html=True)
        cr = fo.groupby("customer_name")["total_amount"].sum().reset_index().sort_values("total_amount", ascending=False).head(10)
        if not cr.empty:
            cr["label"] = cr["customer_name"].str[:22]
            fig2 = go.Figure(go.Pie(
                labels=cr["label"], values=cr["total_amount"], hole=0.55,
                marker=dict(colors=ACCENT[:len(cr)], line=dict(color="#0B0F1A", width=2)),
                textinfo="percent", textfont=dict(size=11, color="#E5E7EB"),
                hovertemplate="<b>%{label}</b><br>₹%{value:,.0f}<br>%{percent}<extra></extra>",
            ))
            fig2.add_annotation(text="Top 10<br>Customers", x=0.5, y=0.5,
                                font=dict(size=12, color="#9CA3AF"), showarrow=False)
            fig2.update_layout(**CHART_LAYOUT, title="Customer Revenue Distribution", height=320)
            chart(fig2, 320)

    # Order value distribution
    st.markdown(section_header("Order Value Distribution"), unsafe_allow_html=True)
    valid_o = fo[fo["total_amount"] > 0]
    if not valid_o.empty:
        fig3 = go.Figure(go.Histogram(
            x=valid_o["total_amount"], nbinsx=25,
            marker=dict(color="#7C3AED", opacity=0.75, line=dict(color="#A78BFA", width=0.5)),
            hovertemplate="Range: ₹%{x:,.0f}<br>Count: %{y}<extra></extra>",
        ))
        fig3.update_layout(**CHART_LAYOUT, title="Distribution of Order Values", height=260)
        fig3.update_layout(xaxis=dict(tickformat=",.0f", tickprefix="₹", **CHART_LAYOUT["xaxis"]))
        chart(fig3, 260)

    # Detailed order table
    st.markdown(section_header("Order Detail"), unsafe_allow_html=True)
    show_df = fo[["po_number","document_date","vendor_name","customer_name",
                  "consignee_state","total_amount","gst_amount","financial_year"]].copy()
    show_df.columns = ["PO#","Date","Vendor","Customer","State","Amount","GST","FY"]
    show_df = show_df.sort_values("Date", ascending=False)
    show_df["Amount"] = show_df["Amount"].apply(lambda x: fmt(x))
    show_df["GST"]    = show_df["GST"].apply(lambda x: fmt(x) if x else "—")
    st.dataframe(show_df, use_container_width=True, height=360,
                 column_config={
                     "PO#": st.column_config.TextColumn("PO Number", width="medium"),
                     "Date": st.column_config.DateColumn("Date"),
                     "Vendor": st.column_config.TextColumn("Vendor", width="large"),
                     "Customer": st.column_config.TextColumn("Customer", width="large"),
                 })

    buf = io.BytesIO()
    fo.to_excel(buf, index=False, engine="xlsxwriter")
    st.download_button("📥 Export to Excel", buf.getvalue(),
                       "ms_enterprises_financial.xlsx",
                       "application/vnd.ms-excel")



# ═══════════════════════════════════════════════════════════
# PAGE: DATA ENTRY
# ═══════════════════════════════════════════════════════════
elif page == "📝  Data Entry":
    st.markdown('<div class="page-title">Data <span>Entry</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="page-sub">Add missing orders manually · all entries reflect instantly in every chart</div>', unsafe_allow_html=True)

    if True:
        st.markdown(section_header("Add New Order / Missing Entry",
            "Fill details below — saved to database, reflects in all pages"), unsafe_allow_html=True)

        with st.form("manual_po_form", clear_on_submit=True):
            fc1, fc2, fc3 = st.columns(3)
            with fc1:
                f_po      = st.text_input("PO Number *", placeholder="e.g. PO-2024-001")
                f_date    = st.date_input("Order Date *")
                f_fy      = st.selectbox("Financial Year *",
                                         ["2021-22","2022-23","2023-24","2024-25","2025-26"],
                                         index=3)
            with fc2:
                f_client  = st.text_input("Factory / Client Name *", placeholder="e.g. Adani Power Ltd")
                f_state   = st.selectbox("State *", sorted([
                    "Maharashtra","Chhattisgarh","Jharkhand","Madhya Pradesh","Odisha",
                    "Rajasthan","West Bengal","Andhra Pradesh","Uttar Pradesh","Telangana",
                    "Karnataka","Tamil Nadu","Gujarat","Delhi","Bihar","Haryana","Punjab","Other"
                ]))
                f_vendor  = st.text_input("Vendor / Supplier", value="MS ENTERPRISES")
            with fc3:
                f_mat     = st.text_input("Material Name *", placeholder="e.g. Plastic Refractory")
                f_cat     = st.selectbox("Material Category",
                                         sorted(items_df["material_category"].dropna().unique().tolist()
                                                or ["Refractory","Castable","Binder","Other"]))
                f_qty     = st.number_input("Quantity (MT) *", min_value=0.0, step=0.5, format="%.2f")

            fc4, fc5, fc6 = st.columns(3)
            with fc4:
                f_rate    = st.number_input("Rate per MT (₹) *", min_value=0.0, step=100.0, format="%.0f")
            with fc5:
                f_amount  = st.number_input("Total Amount (₹)", min_value=0.0, step=1000.0, format="%.0f",
                                             help="Leave 0 to auto-calculate from Qty × Rate")
            with fc6:
                f_gst     = st.number_input("GST Amount (₹)", min_value=0.0, step=100.0, format="%.0f")

            f_notes = st.text_area("Notes / Remarks", placeholder="Optional remarks…", height=68)
            submitted = st.form_submit_button("💾  Save Entry", use_container_width=True)

        if submitted:
            if not f_po or not f_client or not f_mat or f_qty <= 0 or f_rate <= 0:
                st.error("Fill all required fields (*) and ensure Qty & Rate > 0.")
            else:
                total_amt = f_amount if f_amount > 0 else f_qty * f_rate
                doc_date  = f_date.strftime("%Y-%m-%d")
                month_num = f_date.month
                year_num  = f_date.year

                conn = sqlite3.connect(DB_PATH)
                cur  = conn.cursor()

                # Upsert material
                cur.execute("SELECT id FROM materials WHERE name_standardized=?", (f_mat,))
                mat_row = cur.fetchone()
                if mat_row:
                    mat_id = mat_row[0]
                else:
                    cur.execute("INSERT INTO materials (name_standardized, category) VALUES (?,?)",
                                (f_mat, f_cat))
                    mat_id = cur.lastrowid

                # Insert order
                cur.execute("""
                    INSERT INTO orders
                      (po_number, document_date, consignee_name, consignee_state,
                       supplier_name, total_amount, gst_amount,
                       financial_year, month, year, business_type)
                    VALUES (?,?,?,?,?,?,?,?,?,?,?)
                """, (f_po, doc_date, f_client, f_state, f_vendor,
                      total_amt, f_gst, f_fy, month_num, year_num,
                      f"MANUAL | {f_notes}" if f_notes else "MANUAL"))
                order_id = cur.lastrowid

                # Insert order item
                cur.execute("""
                    INSERT INTO order_items
                      (order_id, material_id, quantity, rate, amount, description)
                    VALUES (?,?,?,?,?,?)
                """, (order_id, mat_id, f_qty, f_rate, f_qty * f_rate, f_mat))

                conn.commit()
                conn.close()

                st.success(f"✅ Saved! PO **{f_po}** — {f_mat} · {f_qty} MT · ₹{total_amt:,.0f}")
                st.cache_data.clear()
                st.rerun()

        # Show recent manual entries (orders not from PDF processing)
        st.markdown(section_header("Recent Entries"), unsafe_allow_html=True)
        try:
            conn = sqlite3.connect(DB_PATH)
            recent = pd.read_sql_query("""
                SELECT o.po_number, o.document_date, o.consignee_name, o.consignee_state,
                       o.total_amount, o.financial_year,
                       GROUP_CONCAT(m.name_standardized, ', ') as materials,
                       SUM(oi.quantity) as total_qty
                FROM orders o
                LEFT JOIN order_items oi ON oi.order_id = o.id
                LEFT JOIN materials m   ON m.id = oi.material_id
                WHERE o.business_type LIKE 'MANUAL%'
                GROUP BY o.id
                ORDER BY o.id DESC LIMIT 30
            """, conn)
            conn.close()
            if not recent.empty:
                recent["total_amount"] = recent["total_amount"].apply(fmt)
                recent["total_qty"]    = recent["total_qty"].apply(lambda x: f"{x:,.1f} MT" if x else "—")
                recent.columns = ["PO#","Date","Client","State","Amount","FY","Materials","Qty"]
                st.dataframe(recent, use_container_width=True, hide_index=True, height=300)
        except Exception as e:
            st.info("No manual entries yet.")

    # ── FULL ORDER TABLE ────────────────────────────────────
    st.markdown("---")
    st.markdown(section_header("All Orders — Full Detail",
        "Every PO with invoice, client, material, rate · use filters to narrow down"),
        unsafe_allow_html=True)

    try:
        conn_t = sqlite3.connect(DB_PATH)
        full_tbl = pd.read_sql_query("""
            SELECT
                o.po_number          AS "Invoice / PO#",
                o.document_date      AS "Date",
                o.financial_year     AS "FY",
                o.customer_name      AS "Client",
                o.consignee_state    AS "State",
                m.name_standardized  AS "Material",
                m.category           AS "Category",
                oi.quantity          AS "Qty (MT)",
                oi.rate              AS "Rate (₹/MT)",
                oi.amount            AS "Amount (₹)",
                o.total_amount       AS "PO Total (₹)"
            FROM orders o
            LEFT JOIN order_items oi ON oi.order_id = o.id
            LEFT JOIN materials   m  ON m.id = oi.material_id
            ORDER BY o.document_date DESC, o.id DESC
        """, conn_t)
        conn_t.close()

        if not full_tbl.empty:
            # KG→MT normalisation for display
            # filter controls
            fc1, fc2, fc3, fc4 = st.columns(4)
            fy_f   = fc1.selectbox("Financial Year", ["All"] + sorted(full_tbl["FY"].dropna().unique().tolist()),
                                   label_visibility="collapsed")
            st_f   = fc2.selectbox("State",   ["All"] + sorted(full_tbl["State"].dropna().unique().tolist()),
                                   label_visibility="collapsed")
            cl_f   = fc3.selectbox("Client",  ["All"] + sorted(full_tbl["Client"].dropna().unique().tolist()),
                                   label_visibility="collapsed")
            mat_f  = fc4.selectbox("Material",["All"] + sorted(full_tbl["Material"].dropna().unique().tolist()),
                                   label_visibility="collapsed")

            disp = full_tbl.copy()
            if fy_f  != "All": disp = disp[disp["FY"]       == fy_f]
            if st_f  != "All": disp = disp[disp["State"]    == st_f]
            if cl_f  != "All": disp = disp[disp["Client"]   == cl_f]
            if mat_f != "All": disp = disp[disp["Material"] == mat_f]

            # format numbers
            disp["Qty (MT)"]     = disp["Qty (MT)"].apply(lambda x: f"{x:,.2f}" if pd.notna(x) and x else "—")
            disp["Rate (₹/MT)"]  = disp["Rate (₹/MT)"].apply(lambda x: f"₹{x:,.0f}" if pd.notna(x) and x else "—")
            disp["Amount (₹)"]   = disp["Amount (₹)"].apply(lambda x: f"₹{x:,.0f}" if pd.notna(x) and x else "—")
            disp["PO Total (₹)"] = disp["PO Total (₹)"].apply(lambda x: f"₹{x:,.0f}" if pd.notna(x) and x else "—")

            st.markdown(f"<span style='font-size:12px;color:#64748B;'>{len(disp):,} rows</span>",
                        unsafe_allow_html=True)
            st.dataframe(disp, use_container_width=True, hide_index=True, height=480)
        else:
            st.info("No order data found.")
    except Exception as e:
        st.error(f"Could not load full order table: {e}")


# ═══════════════════════════════════════════════════════════
# PAGE: PROCESSING
# ═══════════════════════════════════════════════════════════
elif page == "⚙️  Processing":
    st.markdown('<div class="page-title">PDF <span>Processing</span></div>', unsafe_allow_html=True)
    st.markdown('<div class="page-sub">Extraction pipeline status and file processing logs</div>', unsafe_allow_html=True)

    success = len(processing_df[processing_df["status"]=="success"])
    errors  = len(processing_df[processing_df["status"]=="error"])
    empty   = len(processing_df[processing_df["status"]=="empty"])
    total   = len(processing_df)

    pc1, pc2, pc3, pc4 = st.columns(4)
    pc1.metric("Total PDFs",          str(total))
    pc2.metric("Processed OK",        str(success), delta=f"{success/total*100:.0f}% success" if total else "")
    pc3.metric("Errors",              str(errors),  delta=f"-{errors}" if errors else "0", delta_color="inverse")
    pc4.metric("Empty / Non-PO",      str(empty))

    c1, c2 = st.columns([1, 2])
    with c1:
        status_counts = processing_df["status"].value_counts().reset_index()
        status_counts.columns = ["status","count"]
        fig = go.Figure(go.Pie(
            labels=status_counts["status"], values=status_counts["count"], hole=0.55,
            marker=dict(colors=["#10B981","#EF4444","#F59E0B","#6B7280"][:len(status_counts)],
                        line=dict(color="#0B0F1A", width=2)),
            textinfo="percent+label", textfont=dict(size=11, color="#E5E7EB"),
        ))
        fig.add_annotation(text=f"{total}<br>Files", x=0.5, y=0.5,
                           font=dict(size=14, color="#E5E7EB"), showarrow=False)
        fig.update_layout(**CHART_LAYOUT, title="Processing Status", height=300)
        fig.update_layout(showlegend=False)
        chart(fig, 300)

    with c2:
        st.markdown(section_header("Processing Log"), unsafe_allow_html=True)
        log_df = processing_df[["file_path","status","pages","ocr_confidence","processed_at","error_message"]].copy()
        log_df["file_path"] = log_df["file_path"].apply(lambda x: os.path.basename(x) if isinstance(x, str) else x)
        log_df.columns = ["File","Status","Pages","Confidence","Processed At","Error"]
        st.dataframe(log_df.sort_values("Processed At", ascending=False),
                     use_container_width=True, height=320)

    if st.button("🔄 Reprocess All PDFs"):
        st.warning("Run from terminal:  `python processor.py --force`  to reprocess all 164 PDFs")
